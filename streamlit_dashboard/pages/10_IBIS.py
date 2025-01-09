import streamlit as st
import pandas as pd
import mysql.connector
from plotly.subplots import make_subplots
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
import io

st.set_page_config(page_title="IBIS-Industry Analysis", layout="wide")

# Function to get distinct industries
def get_industries():
    host = st.secrets["mysql"]["host"]
    user = st.secrets["mysql"]["user"]
    password = st.secrets["mysql"]["password"]
    database = st.secrets["mysql"]["database"]

    # Connect to the database
    connection = mysql.connector.connect(
        host=host,
        user=user,
        password=password,
        database=database
    )

    # Query to get distinct industries
    query = "SELECT DISTINCT Industry FROM ibis_report"
    df = pd.read_sql(query, connection)
    connection.close()
    return df

# Function to get data for the selected industry
def get_data(industry):
    host = st.secrets["mysql"]["host"]
    user = st.secrets["mysql"]["user"]
    password = st.secrets["mysql"]["password"]
    database = st.secrets["mysql"]["database"]

    # Connect to the database
    connection = mysql.connector.connect(
        host=host,
        user=user,
        password=password,
        database=database
    )

    # Query to get data for the selected industry
    query = f"SELECT * FROM ibis_report WHERE Industry = '{industry}'"
    df = pd.read_sql(query, connection)
    connection.close()
    return df

def create_category_charts(df):
    # Initialize charts
    charts = []

    # Define bar and line colors
    bar_color = '#032649'
    line_color = '#EB8928'

    # Loop through the categories
    for category in ['Profit', 'Revenue', 'Business', 'Employees']:
        if category in df['Category'].unique():
            category_data = df[df['Category'] == category]

            # Create a subplot with secondary y-axis
            fig = make_subplots(specs=[[{"secondary_y": True}]])
            
            # Add bar chart for 'Value' on the primary y-axis
            fig.add_trace(
                go.Bar(
                    x=category_data['Year'],
                    y=category_data['Value'],
                    name='Value',
                    marker_color=bar_color,
                    text=[f"{value}" if i == len(category_data) - 1 else "" for i, value in enumerate(category_data['Value'])],
                    textposition="outside"
                ),
                secondary_y=False
            )

            # Add line chart for 'Change' on the secondary y-axis
            fig.add_trace(
                go.Scatter(
                    x=category_data['Year'],
                    y=category_data['Change'],
                    name='Change (%)',
                    mode='lines+markers',
                    line=dict(color=line_color),
                    text=[f"{change:.1f}%" if i == len(category_data) - 1 else "" for i, change in enumerate(category_data['Change'])],
                    textposition="top center"
                ),
                secondary_y=True
            )

            # Update layout to include category name as title
            fig.update_layout(
                title=dict(
                    text=f"{category}",
                    font=dict(size=16, color="#595959"),
                    x=0.5,  # Center-align title
                    xanchor='center'
                ),
                xaxis_title="Year",
                yaxis_title="Value",
                legend=dict(x=0, y=1, xanchor='left', yanchor='top'),
                xaxis=dict(showgrid=False, color="#595959",  
                    tickfont=dict(color="#595959")),
                yaxis=dict(showgrid=False, color="#595959",
                    tickfont=dict(color="#595959")),
                margin=dict(l=20, r=20, t=50, b=50),
                height=400,
                width=600
            )
            fig.update_yaxes(title_text="Value (in bn$)", secondary_y=False)
            fig.update_yaxes(title_text="Change (%)", secondary_y=True)

            # Append the figure to the list of charts
            charts.append(fig)

    return charts

# Function to export charts to PowerPoint
def export_charts_to_ppt(charts, filename="charts.pptx"):
    prs = Presentation()
    
    # Define chart positions (in Inches)
    positions = [
        {"left": Inches(1), "top": Inches(1), "width": Inches(6), "height": Inches(3)},  # Position for Chart 1
        {"left": Inches(1), "top": Inches(4.5), "width": Inches(6), "height": Inches(3)},  # Position for Chart 2
        {"left": Inches(7.5), "top": Inches(1), "width": Inches(6), "height": Inches(3)},  # Position for Chart 3
        {"left": Inches(7.5), "top": Inches(4.5), "width": Inches(6), "height": Inches(3)},  # Position for Chart 4
    ]

    for i, chart in enumerate(charts):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout

        image_stream = io.BytesIO()
        chart.write_image(image_stream, format='png')
        image_stream.seek(0)

        pos = positions[i % len(positions)]

        slide.shapes.add_picture(
            image_stream,
            pos["left"],
            pos["top"],
            width=pos["width"],
            height=pos["height"]
        )

    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

st.title("IBIS - Industry Analysis")

df_industries = get_industries()
industry_options = df_industries["Industry"].tolist()
industry = st.selectbox("Select Industry", industry_options)

# Get the data for the selected industry
df = get_data(industry)

# Create separate charts for each category
if not df.empty:
    charts = create_category_charts(df)
    
    # Display all charts
    for chart in charts:
        st.plotly_chart(chart)

    # Display debug message
    st.write("Charts are ready for export!")

    # Add "Export to PowerPoint" button
    if st.button("Export Charts to PowerPoint"):
        ppt_buffer = export_charts_to_ppt(charts)

        # Provide download link for PowerPoint
        st.download_button(
            label="Download PowerPoint",
            data=ppt_buffer,
            file_name="industry_charts.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
else:
    st.write("No data available for the selected industry.")
