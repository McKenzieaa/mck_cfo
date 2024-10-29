import streamlit as st
import pandas as pd
import requests
import zipfile
import io
from datetime import date
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import seaborn as sns
import matplotlib.pyplot as plt

# Streamlit page configuration
st.set_page_config(
    page_title="State Indicators Dashboard",
    layout="wide"
)

# Initialize data
today_date = date.today().strftime("%Y-%m-%d")
state_gdp_data = None  # Initialize state GDP data

# State Data IDs
states_data_id = {
    "Alabama": {"ur_id": "ALUR", "labour_id": "LBSSA01"},
    "Alaska": {"ur_id": "AKUR", "labour_id": "LBSSA02"},
    "Arizona": {"ur_id": "AZUR", "labour_id": "LBSSA04"},
    # Add remaining states...
    "Wyoming": {"ur_id": "WYUR", "labour_id": "LBSSA56"}
}

# Set seaborn and matplotlib configurations
sns.set_style("white")
plt.rc('axes', edgecolor='white')
plt.rc('xtick', color='gray')
plt.rc('ytick', color='gray')

def download_csv(state_name, data_type):
    data_ids = states_data_id.get(state_name)
    if not data_ids:
        return None

    data_id = data_ids["ur_id"] if data_type == "unemployment" else data_ids["labour_id"]
    url = f"https://fred.stlouisfed.org/graph/fredgraph.csv?id={data_id}&cosd=1976-01-01&coed={today_date}"

    response = requests.get(url)
    if response.status_code == 200:
        csv_data = pd.read_csv(io.StringIO(response.content.decode("utf-8")))
        column_name = "Unemployment" if data_type == "unemployment" else "Labour Force"
        csv_data.rename(columns={csv_data.columns[1]: column_name}, inplace=True)
        csv_data['DATE'] = pd.to_datetime(csv_data['DATE'])
        return csv_data
    else:
        st.error(f"Error downloading {data_type} data for {state_name}.")
        return None

def load_state_gdp_data():
    """Download and preprocess state-level GDP data."""
    global state_gdp_data
    url = "https://apps.bea.gov/regional/zip/SAGDP.zip"
    response = requests.get(url)

    if response.status_code == 200:
        with zipfile.ZipFile(io.BytesIO(response.content)) as z:
            csv_file_name = next((name for name in z.namelist() 
                                  if name.startswith("SAGDP1__ALL_AREAS_") and name.endswith(".csv")), None)

            if csv_file_name:
                with z.open(csv_file_name) as f:
                    df = pd.read_csv(f, usecols=lambda col: col not in [
                        "GeoFIPS", "Region", "TableName", "LineCode", 
                        "IndustryClassification", "Unit"
                    ], dtype={"Description": str})
                df = df[df["Description"] == "Current-dollar GDP (millions of current dollars) "]
                df = df.melt(id_vars=["GeoName"], var_name="Year", value_name="Value")
                df.rename(columns={"GeoName": "State"}, inplace=True)
                df = df[df["Year"].str.isdigit()]
                df["Year"] = df["Year"].astype(int)
                state_gdp_data = df[df["State"] != "United States"]
            else:
                st.error("No CSV file found with the specified prefix.")
    else:
        st.error("Failed to download GDP data.")
# Load GDP data on startup
load_state_gdp_data()

def plot_unemployment_labour_chart(state_name):
    """Plot unemployment and labour force trends."""
    unemployment_data = download_csv(state_name, "unemployment")
    labour_data = download_csv(state_name, "labour")

    if unemployment_data is not None and labour_data is not None:
        # Filter data to start from the year 2000
        unemployment_data = unemployment_data[unemployment_data['DATE'].dt.year >= 2000]
        labour_data = labour_data[labour_data['DATE'].dt.year >= 2000]

        merged_data = pd.merge(unemployment_data, labour_data, on='DATE')

        # Set figure size inside the function
        fig, ax = plt.subplots(figsize=(6, 4))
        sns.lineplot(data=merged_data, x='DATE', y='Unemployment', 
                     color='#032649', label='Unemployment', ax=ax)
        sns.lineplot(data=merged_data, x='DATE', y='Labour Force', 
                     color='#EB8928', label='Labour Force', ax=ax)

        # Customize axes
        ax.tick_params(axis='both', colors='gray')  # Set tick color to gray
        ax.set_xlabel('')  # Clear x-axis label
        ax.set_ylabel("Value", color='gray')  # Set y-axis label with gray color

        # Add labels for 50% of the data points
        # for i, row in merged_data.iloc[::30].iterrows():  # Label every other data point
        #     ax.text(row['DATE'], row['Unemployment'], 
        #             f"{row['Unemployment']:.1f}", color='black', ha='right', fontsize=8)
        #     ax.text(row['DATE'], row['Labour Force'], 
        #             f"{row['Labour Force']:.1f}", color='black', ha='right', fontsize=8)
        last_row = merged_data.iloc[-1]
        ax.text(last_row['DATE'], last_row['Unemployment'], 
                f"{last_row['Unemployment']:.1f}", color='#032649', ha='right')
        ax.text(last_row['DATE'], last_row['Labour Force'], 
                f"{last_row['Labour Force']:.1f}", color='#EB8928', ha='right')

        ax.legend(loc='center right',fontsize=8, frameon=False)
        st.pyplot(fig)
        return fig
    else:
        st.warning(f"No data available for {state_name}.")
        return None

def plot_gdp_chart(state_name):
    """Plot GDP trends over time."""
    global state_gdp_data

    if state_gdp_data is not None:
        # Filter data to start from the year 2000
        gdp_data = state_gdp_data[state_gdp_data["State"].str.lower() == state_name.lower()]
        gdp_data = gdp_data[gdp_data["Year"] >= 2000]

        if not gdp_data.empty:
            # Set figure size inside the function
            fig, ax = plt.subplots(figsize=(6, 4))
            sns.lineplot(data=gdp_data, x='Year', y='Value', color='black', label='GDP', ax=ax)

            # Customize axes
            ax.tick_params(axis='both', colors='gray')  # Set tick color to gray
            ax.set_xlabel('')  # Clear x-axis label
            ax.set_ylabel("GDP (Millions)", color='gray')  # Set y-axis label

            # # Add labels for 50% of the data points
            # for i, row in gdp_data.iloc[::25].iterrows():  # Label every other data point
            #     ax.text(row['Year'], row['Value'], 
            #             f"{row['Value']:.1f}", color='#032649', ha='right', fontsize=10)
            last_row = gdp_data.iloc[-1]
            ax.text(last_row['Year'], last_row['Value'], 
                    f"{last_row['Value']:.1f}", color='#032649', ha='right')
            ax.legend(loc='upper left',fontsize=8, frameon=False)
            st.pyplot(fig)
            return fig
        else:
            st.warning(f"No GDP data available for {state_name}.")
            st.write("Available State Names in GDP Data:", state_gdp_data["State"].unique())
            return None
    else:
        st.warning("State GDP data not loaded.")
        return None
    
def export_to_pptx(labour_fig, gdp_fig):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]

    slide1 = prs.slides.add_slide(slide_layout)
    title1 = slide1.shapes.title
    title1.text = "Unemployment & Labour Force"
    img1 = BytesIO()
    labour_fig.savefig(img1, format="png")
    img1.seek(0)
    slide1.shapes.add_picture(img1, Inches(1), Inches(1), width=Inches(10))

    slide2 = prs.slides.add_slide(slide_layout)
    title2 = slide2.shapes.title
    title2.text = "GDP"
    img2 = BytesIO()
    gdp_fig.savefig(img2, format="png")
    img2.seek(0)
    slide2.shapes.add_picture(img2, Inches(1), Inches(1), width=Inches(10))

    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def get_state_indicators_layout():
    """Render the State Indicators page layout with side-by-side charts."""
    state_name = st.selectbox("Select State", list(states_data_id.keys()), index=0)

    # Create two columns for side-by-side charts
    col1, col2 = st.columns(2)

    # Plot unemployment and labour chart in the first column
    with col1:
        st.write(f"### {state_name} - Unemployment & Labour Force")
        labour_fig = plot_unemployment_labour_chart(state_name)  # Plot with default figure size inside the function

    # Plot GDP chart in the second column
    with col2:
        st.write(f"### {state_name} - GDP Over Time")
        gdp_fig = plot_gdp_chart(state_name)  # Plot with default figure size inside the function

    # Export to PowerPoint button
    if st.button("Export Charts to PowerPoint") and labour_fig and gdp_fig:
        pptx_file = export_to_pptx(labour_fig, gdp_fig)
        st.download_button(
            label="Download PowerPoint",
            data=pptx_file,
            file_name="state_indicators.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
# Render the layout
get_state_indicators_layout()
