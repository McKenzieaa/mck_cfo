import requests
import numpy as np
import zipfile
import io
import dask.dataframe as dd
import plotly.graph_objs as go
import streamlit as st
import plotly.express as px
import pandas as pd
from st_aggrid import AgGrid, GridOptionsBuilder, GridUpdateMode
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from io import BytesIO
from plotly.subplots import make_subplots
from sqlalchemy import create_engine
import os
from datetime import date
import mysql.connector

today = date.today().strftime("%Y-%m-%d")
state_gdp_data = None

host = st.secrets["mysql"]["host"]
user = st.secrets["mysql"]["user"]
password = st.secrets["mysql"]["password"]
database = st.secrets["mysql"]["database"]

engine = create_engine(f"mysql+pymysql://{user}:{password}@{host}/{database}")

def fetch_data(query):
    try:
        with engine.connect() as connection:
            df = pd.read_sql(query, connection)
        return df
    except Exception as e:
        st.error(f"Error fetching data: {e}")
        return pd.DataFrame()

try:
    conn = mysql.connector.connect(
        host=host,
        user=user,
        password=password,
        database=database
    )
except mysql.connector.Error as e:
    st.error(f"Error connecting to MySQL: {e}")
    st.stop()

query1 = """
SELECT 
    `Year`, `Target`, `EV/Revenue`, `EV/EBITDA`, `Business Description`, `Industry`, `Location`
FROM 
    precedent_table
"""
try:
    df_precedent = pd.read_sql(query1, conn)
except Exception as e:
    st.error(f"Error loading data from MySQL (Precedent Transactions): {e}")
    st.stop()

# Fetch Public Listed Companies data
query2 = """
SELECT `Name`, `Country`, `Enterprise Value (in $)`, `Revenue (in $)`, `EBITDA (in $)`, `Business Description`, `Industry`
FROM public_comp_table
"""
try:
    df_public = pd.read_sql(query2, conn)
except Exception as e:
    st.error(f"Error loading data from MySQL (Public Companies): {e}")
    st.stop()

query3 = """
SELECT `NAICS`, `LineItems`, `Percent`, `ReportID`, `Industry`
FROM rma_table
"""
try:
    df_rma = pd.read_sql(query3, conn)
except Exception as e:
    st.error(f"Error loading data from MySQL (Public Companies): {e}")
    st.stop()

conn.close()

try:
    storage_options = {
        'key': st.secrets["aws"]["AWS_ACCESS_KEY_ID"],
        'secret': st.secrets["aws"]["AWS_SECRET_ACCESS_KEY"],
        'client_kwargs': {'region_name': st.secrets["aws"]["AWS_DEFAULT_REGION"]}
    }
except KeyError:
    st.error("AWS credentials are not configured correctly in Streamlit secrets.")
    st.stop()

def update_figure_slide(ppt, title, fig, slide_number, width, height, left, top):
    if fig is None:
        print(f"Skipping slide '{title}' because the figure is None.")
        return  # Skip if fig is None

    slide = ppt.slides[slide_number] 

    fig_image = BytesIO()
    fig.write_image(fig_image, format="png") 
    fig_image.seek(0)

    # Use Inches for size and position only in the add_picture() function
    slide.shapes.add_picture(fig_image, Inches(left), Inches(top), Inches(width), Inches(height))
    fig_image.close()

def add_table_to_slide(slide, df, left, top, width, height, font_size=Pt(10), header_font_size=Pt(12)):
    # Create a table shape on the slide
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(left), Inches(top), Inches(width), Inches(height))

    # Style the header row
    for col_num, col_name in enumerate(df.columns):
        cell = table.table.cell(0, col_num)
        cell.text = str(col_name)
        # Set header font style
        cell.text_frame.paragraphs[0].font.size = header_font_size
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font for header

    # Style the data rows
    for row_num, row in enumerate(df.values):
        for col_num, value in enumerate(row):
            cell = table.table.cell(row_num + 1, col_num)
            cell.text = str(value)
            # Set data cell font style
            cell.text_frame.paragraphs[0].font.size = font_size
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font for data

            # Optional: Adjust vertical alignment and wrapping
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = True

    # Optional: Adjust cell padding (top, bottom, left, right)
    for row in table.table.rows:
        for cell in row.cells:
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)

def export_all_to_pptx(
    labour_fig_us, external_fig, gdp_fig_us, cpi_ppi_fig,
    fig1_precedent, fig2_precedent, fig1_public, fig2_public,
    labour_fig, gdp_fig, income_statement_df, 
    balance_sheet_df, state_name, fig1_ibis, fig2_ibis, fig3_ibis, fig4_ibis
):
    # Load the custom template
    template_path = os.path.join(os.getcwd(), "streamlit_dashboard", "data", "main_template_pitch.pptx")
    ppt = Presentation(template_path)  # Load the template

    # Use the existing slides (slide_number corresponds to the slide index)
    update_figure_slide(ppt, "Precedent - EV/Revenue", fig1_precedent, slide_number=13, width=9, height=3, left=0.45, top=0.90)
    update_figure_slide(ppt, "Precedent - EV/EBITDA", fig2_precedent, slide_number=13, width=9, height=3, left=0.45, top=3.60)
    update_figure_slide(ppt, "Public Comps - EV/Revenue", fig1_public, slide_number=12, width=9, height=3, left=0.45, top=0.90)
    update_figure_slide(ppt, "Public Comps - EV/EBITDA", fig2_public, slide_number=12, width=9, height=3, left=0.45, top=3.60)
    update_figure_slide(ppt, "Labour Force & Unemployment", labour_fig_us, slide_number=5, width=4.8, height=2.50, left=0.08, top=1.3)
    update_figure_slide(ppt, "External Driver Indicators", external_fig, slide_number=7, width=8.5, height=2.5, left=0.20, top=4.35)
    update_figure_slide(ppt, "GDP by Industry", gdp_fig_us, slide_number=5, width=4.5, height=2.50, left=0.08, top=4.4)
    update_figure_slide(ppt, "CPI and PPI Comparison", cpi_ppi_fig, slide_number=4, width=4.5, height=2.50, left=5.10, top=1.3)
    update_figure_slide(ppt, f"Labour force Statistics {state_name}", labour_fig, slide_number=4, width=4.5, height=2.50, left=0.08, top=1.3)
    update_figure_slide(ppt, f"GDP - {state_name}", gdp_fig, slide_number=4, width=4.5, height=2.50, left=0.08, top=4.4)

    # Add IBIS charts
    update_figure_slide(ppt, "Profit - IBIS", fig1_ibis, slide_number=8, width=4.75, height=2.45, left=0.18, top=4.34)
    update_figure_slide(ppt, "Revenue - IBIS", fig2_ibis, slide_number=8, width=4.75, height=2.45, left=0.18, top=1.3)
    update_figure_slide(ppt, "Business - IBIS", fig3_ibis, slide_number=8, width=4.75, height=2.45, left=5.2, top=1.3)
    update_figure_slide(ppt, "Employees - IBIS", fig4_ibis, slide_number=8, width=4.75, height=2.45, left=5.2, top=4.34)

    # Add Benchmarking Tables to Slide
    slide = ppt.slides[9]
    add_table_to_slide(slide, income_statement_df, left=0.35, top=0.90, width=4.3, height=3.4, header_font_size=Pt(12))
    add_table_to_slide(slide, balance_sheet_df, left=5.2, top=0.9, width=4.3, height=5.65, header_font_size=Pt(12))

    # Save the PPT file to BytesIO and return the bytes
    ppt_bytes = BytesIO()
    ppt.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes

# Streamlit page configuration
st.set_page_config(page_title="Pitch Book", layout="wide")

# Define URLs and Paths
country = "USA"
timefrom = 2010
LABOR_QUERY = "SELECT * FROM industry_db.labor_market_data"
df_lfs = fetch_data(LABOR_QUERY)
df_lfs = df_lfs.rename(columns={'labor_force_rate': 'labor_force_rate', 'unemployment_rate': 'unemployment_rate', 'population': 'population'})
df_lfs['year'] = df_lfs['year'].astype('Int64')
df_lfs['month'] = df_lfs['month'].astype('Int64')
df_lfs['population'] = df_lfs['population'] / 1000
df_lfs['population'] = df_lfs['population'].round(2)
df_lfs['unemployment_rate'] = df_lfs['unemployment_rate'].round(2)
df_lfs['labor_force_rate'] = df_lfs['labor_force_rate'].round(2)

# External Driver Data
ED_QUERY = "SELECT code, indicator, year, index_value, __change FROM industry_db.external_drivers"
external_driver_df = fetch_data(ED_QUERY)
external_driver_df = external_driver_df.rename(columns={'__change': '% Change', 'year': 'Year', 'indicator': 'Indicator'})

# Load CPI data
CPI_QUERY = "SELECT series_id, industry, year, month, cpi_value FROM industry_db.cpi_ind_db"
df_cpi = fetch_data(CPI_QUERY)
df_cpi = df_cpi.rename(columns={'cpi_value': 'CPI Value', 'year': 'Year', 'month': 'Month', 'industry': 'Industry'})
df_cpi['Month'] = df_cpi['Month'].astype(str)
df_cpi['Month & Year'] = df_cpi['Month'] + " " + df_cpi['Year'].astype(str)
# df_cpi['Month & Year'] = pd.to_datetime(df_cpi['Month & Year'], format='%m %Y', errors='coerce')
df_cpi = df_cpi.dropna(subset=['series_id', 'Month & Year', 'CPI Value'])
all_items_data = df_cpi[df_cpi['Industry'] == 'All items']
cpi_industry = df_cpi['Industry'].unique().tolist()

# Query to fetch PPI Data
PPI_QUERY = "SELECT date, value FROM industry_db.ppi_us_m"
df_ppi = fetch_data(PPI_QUERY)
df_ppi = df_ppi.rename(columns={'date': 'Date', 'value': 'PPI Value'})
df_ppi['Date'] = pd.to_datetime(df_ppi['Date'], errors='coerce')
df_ppi = df_ppi[df_ppi['Date'] >= '2010-01-01']
df_ppi['Month & Year'] = df_ppi['Date'].dt.strftime('%b %Y')
df_ppi = df_ppi.dropna()

# GDP by Industry
GDP_QUERY = "SELECT * FROM industry_db.gdp_industry"
df_gdp_us = fetch_data(GDP_QUERY)
df_gdp_us = df_gdp_us.rename(columns={"industry": "Industry", "year": "Year", "value": "Value"})
df_gdp_us["Value"] = df_gdp_us["Value"].astype(float)
# GDP PCT by Industry
GDP_PCT_QUERY = "SELECT * FROM industry_db.gdp_pct_ind"
df_pct_change = fetch_data(GDP_PCT_QUERY)
df_pct_change = df_pct_change.rename(columns={"industry": "Industry", "year": "Year", "value": "Percent Change"})
df_pct_change["Percent Change"] = df_pct_change["Percent Change"].astype(float)
# Combine the two datasets
df_combined = pd.merge(df_gdp_us, df_pct_change, on=["Industry", "Year"])
df_gdp_filtered = df_combined[df_combined['Industry'] == 'Gross Domestic Product']
gdp_industry = df_combined['Industry'].unique().tolist()
gdp_industry.remove('GDP')

def get_industries():
    host = st.secrets["mysql"]["host"]
    user = st.secrets["mysql"]["user"]
    password = st.secrets["mysql"]["password"]
    database = st.secrets["mysql"]["database"]

        # Connect to the database
    connection = mysql.connector.connect(
        host=host,
        user=user,
        password=password,
        database=database
        )

        # Query to get distinct industries
    query = "SELECT DISTINCT Industry FROM ibis_report"
    df_ibis = pd.read_sql(query, connection)
    connection.close()
    return df_ibis

    # Function to get data for the selected industry

def get_data(industry):
        host = st.secrets["mysql"]["host"]
        user = st.secrets["mysql"]["user"]
        password = st.secrets["mysql"]["password"]
        database = st.secrets["mysql"]["database"]

        # Connect to the database
        connection = mysql.connector.connect(
            host=host,
            user=user,
            password=password,
            database=database
        )

        # Query to get data for the selected industry
        query = f"SELECT * FROM ibis_report WHERE Industry = '{industry}'"
        df = pd.read_sql(query, connection)
        connection.close()
        return df

def create_category_charts(df):

    fig1, fig2, fig3, fig4 = None, None, None, None

    bar_color = '#032649'
    line_color = '#EB8928'

    # Loop through the categories
    for category in ['Profit', 'Revenue', 'Business', 'Employees']:
        if category in df['Category'].unique():
            category_data = df[df['Category'] == category]

            # Set legend labels based on the category
            if category == 'Revenue':
                bar_name = "Total Revenue($)"
                line_name = "Change %"
            elif category == 'Profit':
                bar_name = "Total Profit($)"
                line_name = "Change %"
            elif category == 'Business':
                bar_name = "Total number of businesses"
                line_name = "Change %"
            elif category == 'Employees':
                bar_name = "Total number of employees"
                line_name = "Change %"
            else:
                bar_name = "Value"
                line_name = "Change (%)"

            fig = make_subplots(specs=[[{"secondary_y": True}]])

            fig.add_trace(
                go.Bar(
                    x=category_data['Year'],
                    y=category_data['Value'],
                    name=bar_name,
                    marker_color=bar_color,
                    text=[f"{value}" if i == len(category_data) - 1 else "" for i, value in enumerate(category_data['Value'])],
                    textposition="outside",
                    textfont=dict(size=12)
                ),
                secondary_y=False
            )

            fig.add_trace(
                go.Scatter(
                    x=category_data['Year'],
                    y=category_data['Change'],
                    name=line_name,
                    mode='lines+markers',
                    line=dict(color=line_color),
                    text=[f"{change:.1f}%" if i == len(category_data) - 1 else "" for i, change in enumerate(category_data['Change'])],
                    textposition="top center"
                ),
                secondary_y=True
            )

            # Update layout and axis titles
            fig.update_layout(
                # title=dict(
                #     text=f"{category}",
                #     font=dict(size=16, color="#474747"),
                #     x=0, 
                #     xanchor='left'
                # ),
                xaxis_title=" ",
                yaxis_title=" ",
                legend=dict(x=0, y=1, xanchor='left', yanchor='top', orientation='h'),
                xaxis=dict(showgrid=False, color="#474747",
                           tickfont=dict(color="#474747")),
                yaxis=dict(showgrid=False, color="#474747",
                           tickfont=dict(color="#474747")),
                margin=dict(l=30, r=50, t=30, b=30),
                height=400,
                width=600
            )
            fig.update_yaxes(title_text=" ", secondary_y=False)
            fig.update_yaxes(title_text=" ", secondary_y=True)

            if category == 'Profit':
                fig1 = fig
            elif category == 'Revenue':
                fig2 = fig
            elif category == 'Business':
                fig3 = fig
            elif category == 'Employees':
                fig4 = fig

    return fig1, fig2, fig3, fig4

def plot_labor_unemployment():
    fig = go.Figure()

    # Plot population as an area chart on the primary y-axis
    fig.add_trace(go.Scatter(
        x=pd.to_datetime(df_lfs[['year', 'month']].assign(day=1)),
        y=df_lfs['population'],
        fill='tozeroy',  # Area chart
        fillcolor='#032649',
        name='Population',
        mode='none',
        line=dict(color='#032649'),
        yaxis='y1'
    ))

    # Plot unemployment rate on the secondary y-axis
    fig.add_trace(go.Scatter(
        x=pd.to_datetime(df_lfs[['year', 'month']].assign(day=1)),
        y=df_lfs['unemployment_rate'],
        name='Unemployment Rate',
        mode='lines',
        line=dict(color='#EB8928'),
        yaxis='y2'
    ))

    # Plot labor force participation rate on the secondary y-axis
    fig.add_trace(go.Scatter(
        x=pd.to_datetime(df_lfs[['year', 'month']].assign(day=1)),
        y=df_lfs['labor_force_rate'],
        name='Labor Force Participation Rate',
        mode='lines',
        line=dict(color='#595959'),
        yaxis='y2'
    ))

    fig.update_layout(
        title='',
        xaxis=dict(
            showgrid=False,
            showticklabels=True,
            color="#474747",
            tickfont=dict(color="#474747"),
            tickangle=0,
            automargin=True
        ),
        yaxis=dict(
            showgrid=False,
            title='Population',
            color="#474747",
            tickfont=dict(color="#474747"),
            side='left',
            range=[df_lfs['population'].min(), df_lfs['population'].max() * 1.1],
            tickformat=',',
            ticksuffix='M'
        ),
        yaxis2=dict(
            title='Rate (%)',
            overlaying='y',
            side='right'
        ),
        legend=dict(
            orientation="h",
            x=0.01,
            y=-0.15,
            bgcolor='rgba(255, 255, 255, 0.6)',
            font=dict(size=10)
        ),
        hovermode='x unified',
        template='plotly_white',
        plot_bgcolor='rgba(0,0,0,0)',  # Transparent plot background
        paper_bgcolor='rgba(0,0,0,0)',  # Transparent paper background
        height=300,  # Increased height for better spacing
        width=500,  # Adjusted width for better visualization
        margin=dict(b=50, t=30, l=10, r=10)  # Add more bottom margin for x-axis labels
    )
    st.plotly_chart(fig, use_container_width=True)
    return fig

def plot_external_driver(selected_indicators):
    colors = ['#032649', '#1C798A', '#EB8928', '#595959', '#A5A5A5']

    if not selected_indicators:
        selected_indicators = ["World GDP"]

    fig = go.Figure()

    for i, indicator in enumerate(selected_indicators):
        indicator_data = external_driver_df[external_driver_df['Indicator'] == indicator]

        if '% Change' not in indicator_data.columns:
            raise ValueError(f"Expected '% Change' column not found in {indicator}")

        color = colors[i % len(colors)]

        if isinstance(color, str) and color.startswith('#') and len(color) == 7:
            fig.add_trace(go.Scatter(
                x=indicator_data['Year'],
                y=indicator_data['% Change'],
                mode='lines',
                name=indicator,
                line=dict(color=color),
            ))
        else:
            raise ValueError(f"Invalid color value: {color} for indicator: {indicator}")

    fig.update_layout(
        title=' ',
        xaxis=dict(
            showgrid=False,
            showticklabels=True,
            color="#474747",
            tickfont=dict(color="#474747"),
        ),
        yaxis=dict(
            title='',
            showgrid=False,
            color="#474747",
            tickfont=dict(color="#474747"),
        ),
        hovermode='x',
        legend=dict(
            orientation="h",
            x=0.01,
            y=-0.35,
            xanchor='left',
            yanchor='bottom',
            traceorder='normal',
            font=dict(size=10, color="#474747"),
            bgcolor='rgba(255, 255, 255, 0)',
        ),
        plot_bgcolor='rgba(0,0,0,0)',
        paper_bgcolor='rgba(0,0,0,0)',
        height=250,
        width=850,
        margin=dict(b=60, t=20, l=2, r=10),
    )
    st.plotly_chart(fig, use_container_width=True)
    return fig

def plot_cpi_ppi(selected_cpi_industry):
    fig = go.Figure()
    
        # Filter CPI data for the selected industry
    cpi_data = df_cpi[df_cpi['Industry'] == selected_cpi_industry]
    
    if not cpi_data.empty:
        fig.add_trace(
            go.Scatter(
                x=cpi_data['Month & Year'],
                y=cpi_data['CPI Value'],
                mode='lines',
                name=f'CPI - {selected_cpi_industry}',
                line=dict(color='#032649')
            )
        )
    else:
        st.warning(f"No data available for the selected industry: {selected_cpi_industry}")
    
        # Ensure CPI-US data is correctly formatted and available
    if not all_items_data.empty:
            fig.add_trace(
                go.Scatter(
                    x=all_items_data['Month & Year'],
                    y=all_items_data['CPI Value'],
                    mode='lines',
                    name='CPI-US',
                    line=dict(color='#EB8928', dash='solid')
                )
            )
    else:
            st.warning("No CPI-US All Items data available to display.")
    
        # Ensure PPI data exists and is formatted correctly
    if not df_ppi.empty:
            fig.add_trace(
                go.Scatter(
                    x=df_ppi['Month & Year'],
                    y=df_ppi['PPI Value'],
                    mode='lines',
                    name='PPI-US',
                    line=dict(color='#1C798A')
                )
            )
    else:
            st.warning("No PPI data available to display.")
    
        # Update Layout
    fig.update_layout(
            title='CPI & PPI Comparison',
            xaxis=dict(
                showgrid=False,
                showticklabels=True,
                color="#474747",
                tickfont=dict(color="#474747"),
            ),
            yaxis=dict(
                title='Index Value',
                showgrid=False,
                color="#474747",
                tickfont=dict(color="#474747"),
            ),
            legend=dict(
                orientation="h",
                x=0.01,
                y=-0.3,
                xanchor='left',
                yanchor='bottom',
                bgcolor='rgba(255, 255, 255, 0.6)',
                font=dict(size=10, color="#474747"),
            ),
            hovermode='x unified',
            plot_bgcolor='rgba(0,0,0,0)',
            paper_bgcolor='rgba(0,0,0,0)',
            height=300,
            width=600,
            margin=dict(b=60, t=20, r=15, l=15),
        )
    
    return fig

def plot_gdp_and_industry(selected_industry=None):
    fig = make_subplots(specs=[[{"secondary_y": True}]])

    # 1. Add GDP Value Line (Primary Axis)
    fig.add_trace(
        go.Scatter(
            x=df_gdp_filtered['Year'],
            y=df_gdp_filtered['Value'],
            mode='lines',
            name='GDP - Value',
            fill='tozeroy', 
            fillcolor='#032649', 
            line=dict(color='#032649', width=2),
            marker=dict(size=10)
        ),
        secondary_y=False
    )

    # 2. Add GDP Percent Change Line (Secondary Axis)
    fig.add_trace(
        go.Scatter(
            x=df_gdp_filtered['Year'],
            y=df_gdp_filtered['Percent Change'],
            mode='lines',
            name='GDP-% Change',
            line=dict(color='#A5A5A5', width=2, dash='solid'),
            marker=dict(size=10)
        ),
        secondary_y=True
    )

    # Check if an industry is selected
    if selected_industry:
        df_industry_filtered = df_combined[df_combined['Industry'] == selected_industry]

        # 3. Add Selected Industry Value Line (Primary Axis)
        fig.add_trace(
            go.Scatter(
                x=df_industry_filtered['Year'],
                y=df_industry_filtered['Value'],
                mode='none',
                name=f'GDP-{selected_industry} Value',
                fill='tozeroy',  # Area chart
                fillcolor='#EB8928', 
                line=dict(color='#EB8928', width=2),
                marker=dict(size=10)
            ),
            secondary_y=False
        )

        # 4. Add Selected Industry Percent Change Line (Secondary Axis)
        fig.add_trace(
            go.Scatter(
                x=df_industry_filtered['Year'],
                y=df_industry_filtered['Percent Change'],
                mode='lines',
                name=f'GDP-{selected_industry}(% Change)',
                line=dict(color='#1C798A', width=2, dash='solid'),
                marker=dict(size=10)
            ),
            secondary_y=True
        )

        # Update layout
        fig.update_layout(
            title='',
            xaxis_title='',
            yaxis_title='Value',
            yaxis2_title='% Change',
            xaxis=dict(
                showgrid=False,
                showticklabels=True,
                color="#474747",  
                tickfont=dict(color="#474747"),
            ),
            yaxis=dict(
                title='',
                showgrid=False,
                color="#474747", 
                tickfont=dict(color="#474747"), 
            ),
            legend=dict(
                orientation="h",
                x=0.01, 
                y=-0.5, 
                xanchor='left', 
                yanchor='bottom', 
                bgcolor='rgba(255, 255, 255, 0.6)',
                font=dict(size=10),
                traceorder='normal'
            ),
            template='plotly_white',
            plot_bgcolor='rgba(0,0,0,0)',
            paper_bgcolor='rgba(0,0,0,0)',
            height=450,
            width=700,
            margin=dict(b=120, t=80,l=10, r=10), 
        )
        st.plotly_chart(fig, use_container_width=True)
        return fig

states_data_id = {
    "Alabama": {"ur_id": "ALUR", "labour_id": "LBSSA01"},
    "Alaska": {"ur_id": "AKUR", "labour_id": "LBSSA02"},
    "Arizona": {"ur_id": "AZUR", "labour_id": "LBSSA04"},
    "Arkansas": {"ur_id": "ARUR", "labour_id": "LBSSA05"},
    "California": {"ur_id": "CAUR", "labour_id": "LBSSA06"},
    "Colorado": {"ur_id": "COUR", "labour_id": "LBSSA08"},
    "Connecticut": {"ur_id": "CTUR", "labour_id": "LBSSA09"},
    "Delaware": {"ur_id": "DEUR", "labour_id": "LBSSA10"},
    "District of Columbia": {"ur_id": "DCUR", "labour_id": "LBSSA11"},
    "Florida": {"ur_id": "FLUR", "labour_id": "LBSSA12"},
    "Georgia": {"ur_id": "GAUR", "labour_id": "LBSSA13"},
    "Hawaii": {"ur_id": "HIUR", "labour_id": "LBSSA15"},
    "Idaho": {"ur_id": "IDUR", "labour_id": "LBSSA16"},
    "Illinois": {"ur_id": "ILUR", "labour_id": "LBSSA17"},
    "Indiana": {"ur_id": "INUR", "labour_id": "LBSSA18"},
    "Iowa": {"ur_id": "IAUR", "labour_id": "LBSSA19"},
    "Kansas": {"ur_id": "KSUR", "labour_id": "LBSSA20"},
    "Kentucky": {"ur_id": "KYURN", "labour_id": "LBSSA21"},
    "Louisiana": {"ur_id": "LAUR", "labour_id": "LBSSA22"},
    "Maine": {"ur_id": "MEUR", "labour_id": "LBSSA23"},
    "Maryland": {"ur_id": "MDUR", "labour_id": "LBSSA24"},
    "Massachusetts": {"ur_id": "MAUR", "labour_id": "LBSSA25"},
    "Michigan": {"ur_id": "MIUR", "labour_id": "LBSSA26"},
    "Minnesota": {"ur_id": "MNUR", "labour_id": "LBSSA27"},
    "Mississippi": {"ur_id": "MSUR", "labour_id": "LBSSA28"},
    "Missouri": {"ur_id": "MTUR", "labour_id": "LBSSA29"},
    "Montana": {"ur_id": "MTUR", "labour_id": "LBSSA30"},
    "Nebraska": {"ur_id": "NEUR", "labour_id": "LBSSA31"},
    "Nevada": {"ur_id": "NVUR", "labour_id": "LBSSA32"},
    "New Hampshire": {"ur_id": "NHUR", "labour_id": "LBSSA33"},
    "New Jersey": {"ur_id": "NJURN", "labour_id": "LBSSA34"},
    "New Mexico": {"ur_id": "NMUR", "labour_id": "LBSSA35"},
    "New York": {"ur_id": "NYUR", "labour_id": "LBSSA36"},
    "North Carolina": {"ur_id": "NCUR", "labour_id": "LBSSA37"},
    "North Dakota": {"ur_id": "NDUR", "labour_id": "LBSSA38"},
    "Ohio": {"ur_id": "OHUR", "labour_id": "LBSSA39"},
    "Oklahoma": {"ur_id": "OKUR", "labour_id": "LBSSA40"},
    "Oregon": {"ur_id": "ORUR", "labour_id": "LBSSA41"},
    "Pennsylvania": {"ur_id": "PAUR", "labour_id": "LBSSA42"},
    "Puerto Rico": {"ur_id": "PRUR", "labour_id": "LBSSA43"},
    "Rhode Island": {"ur_id": "RIUR", "labour_id": "LBSSA44"},
    "South Carolina": {"ur_id": "SCUR", "labour_id": "LBSSA45"},
    "South Dakota": {"ur_id": "SDUR", "labour_id": "LBSSA46"},
    "Tennessee": {"ur_id": "TNUR", "labour_id": "LBSSA47"},
    "Texas": {"ur_id": "TXUR", "labour_id": "LBSSA48"},
    "Utah": {"ur_id": "UTUR", "labour_id": "LBSSA49"},
    "Vermont": {"ur_id": "VTUR", "labour_id": "LBSSA50"},
    "Washington": {"ur_id": "WAUR", "labour_id": "LBSSA54"},
    "West Virginia": {"ur_id": "WVUR", "labour_id": "LBSSA53"},
    "Wisconsin": {"ur_id": "WIUR", "labour_id": "LBSSA55"},
    "Wyoming": {"ur_id": "WYUR", "labour_id": "LBSSA56"}
}
line_colors = {
    "unemployment": "#032649",  # Dark blue
    "labour_force": "#EB8928",  # Orange
    "gdp": "#032649", 
}

def download_csv(state_name, data_type):
    data_ids = states_data_id.get(state_name)
    if not data_ids:
        return None

    data_id = data_ids["ur_id"] if data_type == "unemployment" else data_ids["labour_id"]
    url = f"https://fred.stlouisfed.org/graph/fredgraph.csv?id={data_id}&cosd=1976-01-01"

    response = requests.get(url)
    if response.status_code == 200:
        csv_data = pd.read_csv(io.StringIO(response.content.decode("utf-8")))
        column_name = "Unemployment" if data_type == "unemployment" else "Labour Force"
        csv_data.rename(columns={csv_data.columns[1]: column_name}, inplace=True)
        csv_data = csv_data.rename(columns={'observation_date': 'DATE'})
        csv_data['DATE'] = pd.to_datetime(csv_data['DATE'])
        return csv_data
    else:
        st.error(f"Error downloading {data_type} data for {state_name}.")
        return None

def load_state_gdp_data():
    """Download and preprocess state-level GDP data."""
    global state_gdp_data
    url = "https://apps.bea.gov/regional/zip/SAGDP.zip"
    
    try:
        response = requests.get(url)
        if response.status_code == 200:
            with zipfile.ZipFile(io.BytesIO(response.content)) as z:
                csv_file_name = next(
                    (name for name in z.namelist() 
                     if name.startswith("SAGDP1__ALL_AREAS_") and name.endswith(".csv")), 
                    None
                )
                if csv_file_name:
                    with z.open(csv_file_name) as f:
                        df = pd.read_csv( f, 
                            usecols=lambda col: col not in [ "GeoFIPS", "Region", "TableName", "LineCode", "IndustryClassification", "Unit" ],dtype={"Description": str})
                    df = df[df["Description"] == "Current-dollar GDP (millions of current dollars) "]
                    df = df.melt(id_vars=["GeoName"], var_name="Year", value_name="Value")
                    df.rename(columns={"GeoName": "State"}, inplace=True)
                    df = df[df["Year"].str.isdigit()]
                    df["Year"] = df["Year"].astype(int)
                    df["Value"] = pd.to_numeric(df["Value"], errors='coerce')
                    df.dropna(subset=["Value"], inplace=True)
                    state_gdp_data = df[df["State"] != "United States"]
                else:
                    st.error("No matching CSV file found in the downloaded ZIP.")
        else:
            st.error(f"Failed to download GDP data. Status code: {response.status_code}")

    except Exception as e:
        st.error(f"An error occurred: {e}")

load_state_gdp_data()

def plot_unemployment_labour_chart(state_name):
    unemployment_data = download_csv(state_name, "unemployment")
    labour_data = download_csv(state_name, "labour")

    if unemployment_data is not None and labour_data is not None:
        unemployment_data = unemployment_data[unemployment_data['DATE'].dt.year >= 2000]
        labour_data = labour_data[labour_data['DATE'].dt.year >= 2000]

        merged_data = pd.merge(unemployment_data, labour_data, on='DATE')

        fig = go.Figure()
        fig.add_trace(go.Scatter(x=merged_data['DATE'], y=merged_data['Unemployment'], mode='lines',line=dict(color=line_colors["unemployment"]), name="Unemployment"))
        fig.add_trace(go.Scatter(x=merged_data['DATE'], y=merged_data['Labour Force'], mode='lines',line=dict(color=line_colors["labour_force"]), name="Labour Force"))

        last_row = merged_data.iloc[-1]
        fig.add_annotation(
            x=last_row['DATE'], y=last_row['Unemployment'],
            text=f"{last_row['Unemployment']:.1f}"+"%", showarrow=True, arrowhead=1, ax=-40, ay=-40
        )
        fig.add_annotation(
            x=last_row['DATE'], y=last_row['Labour Force'],
            text=f" {last_row['Labour Force']:.1f}"+"%", showarrow=True, arrowhead=1, ax=-40, ay=40
        )

        fig.update_layout(
            title="",
            xaxis_title=" ",
            yaxis_title="Rate",
            template="plotly_white",
            legend=dict(
                x=0.01,  
                y=-0.3,  
                xanchor='left',
                yanchor='bottom',
                title_text=None,
                orientation='h', 
                font=dict(size=10)
            ),
            plot_bgcolor='rgba(0,0,0,0)',
            paper_bgcolor='rgba(0,0,0,0)',
            margin=dict(l=15, r=15, t=10, b=80), 
            height=300,
            width=500,
            xaxis=dict(showgrid=False),
            yaxis=dict(showgrid=False)
        )

        st.plotly_chart(fig, use_container_width=True)
        return fig
    else:
        st.warning(f"No data available for {state_name}.")
        return None

def plot_gdp_chart(state_name):
    global state_gdp_data

    if state_gdp_data is not None:
        gdp_data = state_gdp_data[state_gdp_data["State"].str.lower() == state_name.lower()]
        gdp_data = gdp_data[gdp_data["Year"] >= 2000]

        if not gdp_data.empty:
            fig = go.Figure()
            fig.add_trace(go.Scatter(x=gdp_data['Year'], y=gdp_data['Value'], mode='lines',line=dict(color=line_colors["gdp"]), name=f"{state_name} GDP"))

            last_row = gdp_data.iloc[-1]
            value_in_millions = last_row['Value'] / 1_000_000
            formatted_value = f"{value_in_millions:.1f}M"

            fig.add_annotation(
                x=last_row['Year'], y=last_row['Value'],
                text=f" {formatted_value}", showarrow=True, arrowhead=1, ax=-40, ay=-40
            )

            fig.update_layout(
                title=(""),
                xaxis_title=" ",
                yaxis_title="GDP ($ mn)",
                template="plotly_white",
                legend=dict( x=0.01, y=0.01, xanchor='left', yanchor='bottom',title_text=None ),
                plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',margin=dict(l=2, r=2, t=30,b=50),height=300,width=500,xaxis=dict(showgrid=False, color="#474747",tickfont=dict(color="#474747")),yaxis=dict(showgrid=False, color="#474747",tickfont=dict(color="#474747")))

            st.plotly_chart(fig, use_container_width=True)
            return fig
        else:
            st.warning(f"No GDP data available for {state_name}.")
            return None
    else:
        st.warning("State GDP data not loaded.")
        return None

# Define S3 file paths
# precedent_path = "s3://documentsapi/industry_data/precedent.parquet"
public_comp_path = "s3://documentsapi/industry_data/public_comp_data.parquet"
s3_path_rma = "s3://documentsapi/industry_data/rma_data.parquet"

df_public = df_public.rename(columns={
    'Name': 'Company',
    'Country': 'Location',
    'Enterprise Value (in $)': 'Enterprise Value',
    'Revenue (in $)': 'Revenue',
    'EBITDA (in $)': 'EBITDA',
})
df_public['Enterprise Value'] = pd.to_numeric(df_public['Enterprise Value'], errors='coerce')
df_public['Revenue'] = pd.to_numeric(df_public['Revenue'], errors='coerce')
df_public['EBITDA'] = pd.to_numeric(df_public['EBITDA'], errors='coerce')
df_public['EV/Revenue'] = df_public['Enterprise Value'] / df_public['Revenue']
df_public['EV/EBITDA'] = df_public['Enterprise Value'] / df_public['EBITDA']

precedent_df = df_precedent.copy()
public_comp_df = df_public.copy()

# Accordion for Precedent Transactions
with st.expander("Precedent Transactions"):
    industries = precedent_df['Industry'].dropna().unique()
    locations = precedent_df['Location'].dropna().unique()
    col1, col2 = st.columns(2)
    selected_industries = col1.multiselect("Select Industry", industries, key="precedent_industries")
    selected_locations = col2.multiselect("Select Location", locations, key="precedent_locations")
    if selected_industries and selected_locations:
        filtered_precedent_df = precedent_df[precedent_df['Industry'].isin(selected_industries) & precedent_df['Location'].isin(selected_locations)]
        filtered_precedent_df = filtered_precedent_df[['Target', 'Year', 'EV/Revenue', 'EV/EBITDA','Business Description']]
        filtered_precedent_df['Year'] = filtered_precedent_df['Year'].astype(int)

        st.subheader("Precedent Transactions")
        gb = GridOptionsBuilder.from_dataframe(filtered_precedent_df)
        gb.configure_selection(selection_mode="multiple", use_checkbox=True)
        gb.configure_column(
            field="Target",
            tooltipField="Business Description",
            maxWidth=400
            )
        gb.configure_columns(["Business Description"], hide=False)    
        grid_options = gb.build()

        # Display Ag-Grid table
        grid_response = AgGrid(
            filtered_precedent_df,
            gridOptions=grid_options,
            update_mode=GridUpdateMode.SELECTION_CHANGED,
            height=400,
            width='100%',
            theme='streamlit'
        )
        selected_data = pd.DataFrame(grid_response['selected_rows'])
        if not selected_data.empty:

            avg_data = selected_data.groupby('Year')[['EV/Revenue', 'EV/EBITDA']].mean().reset_index()
            avg_data['Year'] = avg_data['Year'].astype(int)

            # Define colors
            color_ev_revenue = "#032649"  # Default Plotly blue
            color_ev_ebitda = "#032649"   # Default Plotly red

            median_ev_revenue = avg_data['EV/Revenue'].median()
            median_ev_ebitda = avg_data['EV/EBITDA'].median()

            # Create the EV/Revenue chart with data labels
            fig1_precedent = px.bar(avg_data, x='Year', y='EV/Revenue', title="EV/Revenue", text='EV/Revenue')  # No title
            fig1_precedent.update_traces(marker_color=color_ev_revenue, texttemplate='%{text:.1f}'+'x', textposition='auto',textfont=dict(size=12))
            fig1_precedent.update_layout(yaxis_title="EV/Revenue", xaxis_title=" ", bargap=0.4, bargroupgap=0.4, yaxis=dict(showgrid=False),xaxis=dict(tickmode='linear', tick0=avg_data['Year'].min(), dtick=1), shapes=[dict(type='line', x0=avg_data['Year'].min(), x1=avg_data['Year'].max(), y0=median_ev_revenue, y1=median_ev_revenue, line=dict(color='#EB8928', dash='dot', width=2))], annotations=[dict(x=avg_data['Year'].max(), y=median_ev_revenue, xanchor='left', yanchor='bottom', text=f'Median: {median_ev_revenue:.1f}'+'x', showarrow=False, font=dict(size=12, color='gray'), bgcolor='white')],plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',margin=dict(l=0, r=0, t=0),width=900,height=300)

            st.plotly_chart(fig1_precedent)

            # Create the EV/EBITDA chart with data labels
            fig2_precedent = px.bar(avg_data, x='Year', y='EV/EBITDA', title="EV/EBITDA", text='EV/EBITDA')
            fig2_precedent.update_traces(marker_color=color_ev_ebitda, texttemplate='%{text:.1f}'+ 'x', textposition='auto',textfont=dict(size=12))
            fig2_precedent.update_layout(yaxis_title="EV/EBITDA", xaxis_title=" ", bargap=0.4, bargroupgap=0.4, yaxis=dict(showgrid=False),xaxis=dict(tickmode='linear', tick0=avg_data['Year'].min(), dtick=1), shapes=[dict(type='line', x0=avg_data['Year'].min(), x1=avg_data['Year'].max(), y0=median_ev_ebitda, y1=median_ev_ebitda, line=dict(color='#EB8928', dash='dot', width=2))], annotations=[dict(x=avg_data['Year'].max(), y=median_ev_ebitda, xanchor='left', yanchor='bottom', text=f'Median: {median_ev_ebitda:.1f}'+'x', showarrow=False, font=dict(size=12, color='gray'), bgcolor='white')],plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',margin=dict(l=0, r=0, t=0),width=900,height=300)
            
            st.plotly_chart(fig2_precedent)

with st.expander("Public Comps"):
    col1, col2 = st.columns(2)
    industries_public = df_public['Industry'].unique()
    locations_public = df_public['Location'].unique()
    selected_industries = col1.multiselect("Select Industry", industries_public, key="public_industries")
    selected_locations = col2.multiselect("Select Location", locations_public, key="public_locations")
    if selected_industries and selected_locations:
        filtered_df = public_comp_df[public_comp_df['Industry'].isin(selected_industries) & public_comp_df['Location'].isin(selected_locations)]
        filtered_df = filtered_df[['Company',  'EV/Revenue', 'EV/EBITDA', 'Business Description']]
        filtered_df['EV/Revenue'] = filtered_df['EV/Revenue'].round(1)
        filtered_df['EV/EBITDA'] = filtered_df['EV/EBITDA'].round(1)

        # Set up Ag-Grid for selection
        st.subheader("Public Listed Companies")
        gb = GridOptionsBuilder.from_dataframe(filtered_df)
        gb.configure_selection(selection_mode="multiple", use_checkbox=True)
        gb.configure_column(
            field="Company",
            tooltipField="Business Description",
            maxWidth=400
        )
        gb.configure_columns(["Business Description"], hide=False)    
        grid_options = gb.build()

        # Display Ag-Grid table
        grid_response = AgGrid(
            filtered_df,
            gridOptions=grid_options,
            update_mode=GridUpdateMode.SELECTION_CHANGED,
            height=400,
            width='100%',
            theme='streamlit'
        )

        selected_data = pd.DataFrame(grid_response['selected_rows'])

        if not selected_data.empty:
            avg_data = selected_data.groupby('Company')[['EV/Revenue', 'EV/EBITDA']].mean().reset_index()
            avg_data['Company'] = avg_data['Company'].apply(lambda x: '<br>'.join([x[i:i+20] for i in range(0, len(x), 20)]) if len(x) > 20 else x)

            color_ev_revenue = "#032649"  # Default Plotly blue
            color_ev_ebitda = "#032649"   # Default Plotly red

            median_ev_revenue = avg_data['EV/Revenue'].median()
            median_ev_ebitda = avg_data['EV/EBITDA'].median()

            # Create the EV/Revenue chart with data labels
            fig1_public = px.bar(avg_data, x='Company', y='EV/Revenue', title="EV/Revenue", text='EV/Revenue')
            fig1_public.update_traces(marker_color=color_ev_revenue, texttemplate='%{text:.1f}'+'x', textposition='auto',textfont=dict(size=12))
            fig1_public.update_layout(yaxis_title="EV/Revenue", xaxis_title=" ",bargap=0.4,bargroupgap=0.4,yaxis=dict(showgrid=False),xaxis=dict(tickangle=0,automargin=True,tickmode='array',tickvals=avg_data['Company'],ticktext=avg_data['Company']),plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',margin=dict(l=0, r=0, t=50,b=80),width=900,height=300)
            fig1_public.add_shape(type="line",x0=-0.5, x1=len(avg_data['Company']) - 0.5,  y0=median_ev_revenue, y1=median_ev_revenue,line=dict(color="#EB8928", width=2, dash="dot"),  xref="x", yref="y")
            fig1_public.add_annotation(x=len(avg_data['Company']) - 1, y=median_ev_revenue + 0.2, text=f"Median: {median_ev_revenue:.1f}x",showarrow=False, font=dict(size=10, color="gray"), xanchor="left",bgcolor='white')

            st.plotly_chart(fig1_public)

            # Create the EV/EBITDA chart with data labels
            fig2_public = px.bar(avg_data, x='Company', y='EV/EBITDA', title="EV/EBITDA", text='EV/EBITDA')
            fig2_public.update_traces(marker_color=color_ev_ebitda,texttemplate='%{text:.1f}'+'x', textposition='auto',textfont=dict(size=12))
            fig2_public.update_layout(yaxis_title="EV/EBITDA", xaxis_title=" ",bargap=0.4,bargroupgap=0.4,yaxis=dict(showgrid=False),xaxis=dict(tickangle=0,automargin=True,tickmode='array',tickvals=avg_data['Company'],ticktext=avg_data['Company']),plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',margin=dict(l=0, r=0, t=50,b=80),width=900,height=300)
            fig2_public.add_shape(type="line",x0=-0.5, x1=len(avg_data['Company']) - 0.5,  y0=median_ev_ebitda, y1=median_ev_ebitda,line=dict(color="#EB8928", width=2, dash="dot"),  xref="x", yref="y")
            fig2_public.add_annotation(x=len(avg_data['Company']) - 1, y=median_ev_ebitda + 0.2, text=f"Median: {median_ev_ebitda:.1f}x",showarrow=False, font=dict(size=10, color="gray"), xanchor="left",bgcolor='white')
            
            st.plotly_chart(fig2_public)

with st.expander("US Indicators"):
    # st.subheader("US Indicators")

    # Labour Force & Unemployment Data
    st.subheader("Labour Force & Unemployment")
    labour_fig_us = plot_labor_unemployment()

    # External Driver Indicators
    st.subheader("External Driver Indicators")
    selected_indicators = st.multiselect(
        "Select External Indicators",
        options=external_driver_df["Indicator"].unique(),
        default=["World GDP"],
        key="external_indicators_multiselect"
    )
    external_fig = plot_external_driver(selected_indicators)

    # GDP by Industry
    st.subheader("GDP by Industry")
    selected_gdp_industry = st.selectbox(
        "Select Industry",
        options=df_combined["Industry"].unique(),
        index=0,
        key="gdp_industry_selectbox"
    )
    gdp_fig_us = plot_gdp_and_industry(selected_gdp_industry)

    # CPI and PPI Comparison
    st.subheader("CPI & PPI")
    if cpi_industry:
        selected_cpi_industry = st.selectbox("Select CPI Industry", cpi_industry, index=0, key="cpi_industry_selectbox")
        cpi_ppi_fig = plot_cpi_ppi(selected_cpi_industry)
        st.plotly_chart(cpi_ppi_fig, use_container_width=True)
    else:
        st.warning("No industries available for selection.")

with st.expander("State Indicators"):
    st.subheader("State Indicators - US")
    state_name = st.selectbox("Select State", list(states_data_id.keys()), index=0)
    st.subheader(f"{state_name} - Unemployment & Labour Force")
    labour_fig = plot_unemployment_labour_chart(state_name)
    st.subheader(f"{state_name} - GDP Over Time")
    gdp_fig = plot_gdp_chart(state_name)

with st.expander("IBIS"):
    st.subheader("IBIS - Industry Report")

    df_industries = get_industries()  
    industry_options = df_industries["Industry"].tolist()
    industry = st.selectbox("Select Industry", industry_options)

    if industry:
        df_selected = get_data(industry)

        if not df_selected.empty:

            fig1_ibis, fig2_ibis, fig3_ibis, fig4_ibis = create_category_charts(df_selected)

            if fig1_ibis:
                st.subheader("Profit")
                st.plotly_chart(fig1_ibis, use_container_width=True)
            if fig2_ibis:
                st.subheader("Revenue")
                st.plotly_chart(fig2_ibis, use_container_width=True)
            if fig3_ibis:
                st.subheader("Business")
                st.plotly_chart(fig3_ibis, use_container_width=True)
            if fig4_ibis:
                st.subheader("Employees")
                st.plotly_chart(fig4_ibis, use_container_width=True)
        else:
            st.warning(f"No data available for the selected industry: {industry}")

s3_path_public_comp = "s3://documentsapi/industry_data/Public Listed Companies US.xlsx"

# RMA data pre-processing
df_rma = dd.read_parquet(s3_path_rma, storage_options=storage_options)
df_rma = df_rma.rename(columns={
    'ReportID': 'Report_ID',
    'Line Items': 'LineItems',
    'Value': 'Value',
    'Percent': 'Percent'
})

# Load Public Comps data
usecols = [
    "Name", "Industry", "Revenue (in %)", "COGS (in %)", "Gross Profit (in %)", "EBITDA (in %)",
    "Operating Profit (in %)", "Other Expenses (in %)", "Operating Expenses (in %)", "Net Income (in %)",
    "Cash (in %)", "Accounts Receivables (in %)", "Inventories (in %)", "Other Current Assets (in %)",
    "Total Current Assets (in %)", "Fixed Assets (in %)", "PPE (in %)", "Total Assets (in %)",
    "Accounts Payable (in %)", "Short Term Debt (in %)", "Long Term Debt (in %)", "Other Current Liabilities (in %)",
    "Total Current Liabilities (in %)", "Other Liabilities (in %)", "Total Liabilities (in %)",
    "Net Worth (in %)", "Total Liabilities & Equity (in %)"
]

df_public_comp = pd.read_excel(
    s3_path_public_comp,
    sheet_name="FY 2023",
    storage_options=storage_options,
    usecols=usecols,
    engine='openpyxl'
)
df_public_comp = df_public_comp.rename(columns=lambda x: x.replace(" (in %)", ""))

# Unique industries from both sources
industries_rma = df_rma[
    df_rma['Industry'].notnull() & df_rma['Industry'].map(lambda x: isinstance(x, str))
]['Industry'].compute().unique()

industries_public = df_public_comp[
    df_public_comp['Industry'].notnull() & df_public_comp['Industry'].map(lambda x: isinstance(x, str))
]['Industry'].unique()

industries = sorted(set(industries_rma).union(set(industries_public)))

with st.expander("Benchmarking"):
    st.subheader("Benchmarking")
   
    income_statement_items = ["Revenue", "COGS", "Gross Profit", "EBITDA", "Operating Profit", "Other Expenses", "Operating Expenses", "Net Income"]
    balance_sheet_items = ["Cash", "Accounts Receivables", "Inventories", "Other Current Assets", "Total Current Assets", "Fixed Assets", "PPE", "Total Assets", "Accounts Payable", "Short Term Debt", "Long Term Debt", "Other Current Liabilities", "Total Current Liabilities", "Other Liabilities", "Total Liabilities", "Net Worth", "Total Liabilities & Equity"]
    
    selected_industry = st.selectbox("Select Industry", industries)

    if selected_industry:
        # Filter RMA data
        filtered_df_rma = df_rma[df_rma['Industry'] == selected_industry].compute()

        if 'Report_ID' in filtered_df_rma.columns:
            filtered_df_rma['Report_ID'] = filtered_df_rma['Report_ID'].replace({"Assets": "Balance Sheet", "Liabilities & Equity": "Balance Sheet"})

        income_statement_df_rma = filtered_df_rma[filtered_df_rma['Report_ID'] == 'Income Statement'][['LineItems', 'Percent']].rename(columns={'Percent': 'RMA Percent'})
        balance_sheet_df_rma = filtered_df_rma[filtered_df_rma['Report_ID'] == 'Balance Sheet'][['LineItems', 'Percent']].rename(columns={'Percent': 'RMA Percent'})

        # Filter Public Comps data
        filtered_df_public = df_public_comp[df_public_comp['Industry'] == selected_industry]
        df_unpivoted = pd.melt(
            filtered_df_public,
            id_vars=["Name", "Industry"],
            var_name="LineItems",
            value_name="Value"
        )
        df_unpivoted['LineItems'] = df_unpivoted['LineItems'].str.replace(" (in %)", "", regex=False)
        df_unpivoted['Value'] = pd.to_numeric(df_unpivoted['Value'].replace("-", 0), errors='coerce').fillna(0)
        df_unpivoted = df_unpivoted.groupby('LineItems')['Value'].mean().reset_index().rename(columns={'Value': 'Public Comp Percent'})

        # Merge data
        income_statement_df = pd.merge(
            pd.DataFrame({'LineItems': income_statement_items}),
            income_statement_df_rma,
            on='LineItems',
            how='left'
        ).merge(
            df_unpivoted[df_unpivoted['LineItems'].isin(income_statement_items)],
            on='LineItems',
            how='left'
        )

        balance_sheet_df = pd.merge(
            pd.DataFrame({'LineItems': balance_sheet_items}),
            balance_sheet_df_rma,
            on='LineItems',
            how='left'
        ).merge(
            df_unpivoted[df_unpivoted['LineItems'].isin(balance_sheet_items)],
            on='LineItems',
            how='left'
        )

        # Ensure numeric values in required columns by converting to float
        income_statement_df['RMA Percent'] = pd.to_numeric(income_statement_df['RMA Percent'], errors='coerce')
        income_statement_df['Public Comp Percent'] = pd.to_numeric(income_statement_df['Public Comp Percent'], errors='coerce')

        balance_sheet_df['RMA Percent'] = pd.to_numeric(balance_sheet_df['RMA Percent'], errors='coerce')
        balance_sheet_df['Public Comp Percent'] = pd.to_numeric(balance_sheet_df['Public Comp Percent'], errors='coerce')

        # Visualizations
        income_fig = px.bar(
            income_statement_df,
            x="LineItems",
            y=["RMA Percent", "Public Comp Percent"],
            barmode="group",
            text_auto=True
        )
        income_fig.update_layout(xaxis_tickangle=45, height=400)

        balance_fig = px.bar(
            balance_sheet_df,
            x="LineItems",
            y=["RMA Percent", "Public Comp Percent"],
            barmode="group",
            text_auto=True
        )
        balance_fig.update_layout(xaxis_tickangle=45, height=400)

        st.write("Income Statement")
        st.dataframe(income_statement_df, hide_index=True)
        st.plotly_chart(income_fig)

        st.write("Balance Sheet")
        st.dataframe(balance_sheet_df, hide_index=True)
        st.plotly_chart(balance_fig)

if st.button("Export Charts to PowerPoint", key="export_button"):
    try:

        pptx_file = export_all_to_pptx(
            labour_fig_us, external_fig, gdp_fig_us, cpi_ppi_fig, 
            fig1_precedent, fig2_precedent, fig1_public, fig2_public, 
            labour_fig, gdp_fig,
            income_statement_df, balance_sheet_df, state_name,fig1_ibis, fig2_ibis, fig3_ibis, fig4_ibis
        )

        st.download_button(
            label="Download PowerPoint", 
            data=pptx_file,
            file_name=f"Pitch_Book_{date.today().strftime('%Y-%m-%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation" 
        )

    except Exception as e:
        st.error(f"Error during PowerPoint export: {e}")
