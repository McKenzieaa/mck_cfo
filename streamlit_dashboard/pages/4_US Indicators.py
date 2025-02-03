import os
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from PIL import Image
import mysql.connector
from sqlalchemy import create_engine
from datetime import date

today = date.today().strftime("%Y-%m-%d")

# MySQL database connection details
host = st.secrets["mysql"]["host"]
user = st.secrets["mysql"]["user"]
password = st.secrets["mysql"]["password"]
database = st.secrets["mysql"]["database"]
engine = create_engine(f"mysql+pymysql://{user}:{password}@{host}/{database}")

def fetch_data(query):
    try:
        with engine.connect() as connection:
            df = pd.read_sql(query, connection)
        return df
    except Exception as e:
        st.error(f"Error fetching data: {e}")
        return pd.DataFrame()

GDP_QUERY = "SELECT * FROM industry_db.gdp_industry"
GDP_PCT_QUERY = "SELECT * FROM industry_db.gdp_pct_ind"


# labor Force Participation Rate Data
LABOR_QUERY = "SELECT * FROM industry_db.labor_market_data"
df_lfs = fetch_data(LABOR_QUERY)
df_lfs = df_lfs.rename(columns={'labor_force_rate': 'labor_force_rate', 'unemployment_rate': 'unemployment_rate', 'population': 'population'})
df_lfs['year'] = df_lfs['year'].astype('Int64')
df_lfs['month'] = df_lfs['month'].astype('Int64')
df_lfs['population'] = df_lfs['population'] / 1000
df_lfs['population'] = df_lfs['population'].round(2)
df_lfs['unemployment_rate'] = df_lfs['unemployment_rate'].round(2)
df_lfs['labor_force_rate'] = df_lfs['labor_force_rate'].round(2)

# External Driver Data
ED_QUERY = "SELECT code, indicator, year, index_value, __change FROM industry_db.external_drivers"
external_driver_df = fetch_data(ED_QUERY)
external_driver_df = external_driver_df.rename(columns={'__change': '% Change', 'year': 'Year', 'indicator': 'Indicator'})

# Load CPI data
CPI_QUERY = "SELECT series_id, industry, year, month, cpi_value FROM industry_db.cpi_ind_db"
df_cpi = fetch_data(CPI_QUERY)
df_cpi = df_cpi.rename(columns={'cpi_value': 'CPI Value', 'year': 'Year', 'month': 'Month', 'industry': 'Industry'})
df_cpi['Month'] = df_cpi['Month'].astype(str)
df_cpi['Month & Year'] = df_cpi['Month'] + " " + df_cpi['Year'].astype(str)
# df_cpi['Month & Year'] = pd.to_datetime(df_cpi['Month & Year'], format='%m %Y', errors='coerce')
df_cpi = df_cpi.dropna(subset=['series_id', 'Month & Year', 'CPI Value'])
all_items_data = df_cpi[df_cpi['Industry'] == 'All items']
cpi_industry = df_cpi['Industry'].unique().tolist()

# Query to fetch PPI Data
PPI_QUERY = "SELECT date, value FROM industry_db.ppi_us_m"
df_ppi = fetch_data(PPI_QUERY)
df_ppi = df_ppi.rename(columns={'date': 'Date', 'value': 'PPI Value'})
df_ppi['Date'] = pd.to_datetime(df_ppi['Date'], errors='coerce')
df_ppi = df_ppi[df_ppi['Date'] >= '2010-01-01']
df_ppi['Month & Year'] = df_ppi['Date'].dt.strftime('%b %Y')
df_ppi = df_ppi.dropna()

# Fetching and cleaning data
df_gdp_us = fetch_data(GDP_QUERY)
df_gdp_us = df_gdp_us.rename(columns={"industry": "Industry", "year": "Year", "value": "Value"})
df_gdp_us["Value"] = df_gdp_us["Value"].astype(float)

df_pct_change = fetch_data(GDP_PCT_QUERY)
df_pct_change = df_pct_change.rename(columns={"industry": "Industry", "year": "Year", "value": "Percent Change"})
df_pct_change["Percent Change"] = df_pct_change["Percent Change"].astype(float)
# Combine the two datasets
df_combined = pd.merge(df_gdp_us, df_pct_change, on=["Industry", "Year"])
df_gdp_filtered = df_combined[df_combined['Industry'] == 'Gross Domestic Product']
gdp_industry = df_combined['Industry'].unique().tolist()

def plot_labor_unemployment():
    fig = go.Figure()

    # Plot population as an area chart on the primary y-axis
    fig.add_trace(go.Scatter(
        x=pd.to_datetime(df_lfs[['year', 'month']].assign(day=1)),
        y=df_lfs['population'],
        fill='tozeroy',  # Area chart
        fillcolor='#032649',
        name='Population',
        mode='none',
        line=dict(color='#032649'),
        yaxis='y1'
    ))

    # Plot unemployment rate on the secondary y-axis
    fig.add_trace(go.Scatter(
        x=pd.to_datetime(df_lfs[['year', 'month']].assign(day=1)),
        y=df_lfs['unemployment_rate'],
        name='Unemployment Rate',
        mode='lines',
        line=dict(color='#EB8928'),
        yaxis='y2'
    ))

    # Plot labor force participation rate on the secondary y-axis
    fig.add_trace(go.Scatter(
        x=pd.to_datetime(df_lfs[['year', 'month']].assign(day=1)),
        y=df_lfs['labor_force_rate'],
        name='Labor Force Participation Rate',
        mode='lines',
        line=dict(color='#595959'),
        yaxis='y2'
    ))

    fig.update_layout(
        title='',
        xaxis=dict(
            showgrid=False,
            showticklabels=True,
            color="#474747",
            tickfont=dict(color="#474747"),
            tickangle=0,
            automargin=True
        ),
        yaxis=dict(
            showgrid=False,
            title='Population',
            color="#474747",
            tickfont=dict(color="#474747"),
            side='left',
            range=[df_lfs['population'].min(), df_lfs['population'].max() * 1.1],
            tickformat=',',
            ticksuffix='M'
        ),
        yaxis2=dict(
            title='Rate (%)',
            overlaying='y',
            side='right'
        ),
        legend=dict(
            orientation="h",
            x=0.01,
            y=-0.15,
            bgcolor='rgba(255, 255, 255, 0.6)',
            font=dict(size=10)
        ),
        hovermode='x unified',
        template='plotly_white',
        plot_bgcolor='rgba(0,0,0,0)',  # Transparent plot background
        paper_bgcolor='rgba(0,0,0,0)',  # Transparent paper background
        height=300,  # Increased height for better spacing
        width=500,  # Adjusted width for better visualization
        margin=dict(b=50, t=30, l=10, r=10)  # Add more bottom margin for x-axis labels
    )
    st.plotly_chart(fig, use_container_width=True)
    return fig

def plot_external_driver(selected_indicators):
    colors = ['#032649', '#1C798A', '#EB8928', '#595959', '#A5A5A5']

    if not selected_indicators:
        selected_indicators = ["World GDP"]

    fig = go.Figure()

    for i, indicator in enumerate(selected_indicators):
        indicator_data = external_driver_df[external_driver_df['Indicator'] == indicator]

        if '% Change' not in indicator_data.columns:
            raise ValueError(f"Expected '% Change' column not found in {indicator}")

        color = colors[i % len(colors)]

        if isinstance(color, str) and color.startswith('#') and len(color) == 7:
            fig.add_trace(go.Scatter(
                x=indicator_data['Year'],
                y=indicator_data['% Change'],
                mode='lines',
                name=indicator,
                line=dict(color=color),
            ))
        else:
            raise ValueError(f"Invalid color value: {color} for indicator: {indicator}")

    fig.update_layout(
        title=' ',
        xaxis=dict(
            showgrid=False,
            showticklabels=True,
            color="#474747",
            tickfont=dict(color="#474747"),
        ),
        yaxis=dict(
            title='',
            showgrid=False,
            color="#474747",
            tickfont=dict(color="#474747"),
        ),
        hovermode='x',
        legend=dict(
            orientation="h",
            x=0.01,
            y=-0.35,
            xanchor='left',
            yanchor='bottom',
            traceorder='normal',
            font=dict(size=10, color="#474747"),
            bgcolor='rgba(255, 255, 255, 0)',
        ),
        plot_bgcolor='rgba(0,0,0,0)',
        paper_bgcolor='rgba(0,0,0,0)',
        height=250,
        width=850,
        margin=dict(b=60, t=20, l=2, r=10),
    )
    st.plotly_chart(fig, use_container_width=True)
    return fig

def plot_cpi_ppi(selected_cpi_industry):
    fig = go.Figure()
    
        # Filter CPI data for the selected industry
    cpi_data = df_cpi[df_cpi['Industry'] == selected_cpi_industry]
    
    if not cpi_data.empty:
        fig.add_trace(
            go.Scatter(
                x=cpi_data['Month & Year'],
                y=cpi_data['CPI Value'],
                mode='lines',
                name=f'CPI - {selected_cpi_industry}',
                line=dict(color='#032649')
            )
        )
    else:
        st.warning(f"No data available for the selected industry: {selected_cpi_industry}")
    
        # Ensure CPI-US data is correctly formatted and available
    if not all_items_data.empty:
            fig.add_trace(
                go.Scatter(
                    x=all_items_data['Month & Year'],
                    y=all_items_data['CPI Value'],
                    mode='lines',
                    name='CPI-US',
                    line=dict(color='#EB8928', dash='solid')
                )
            )
    else:
            st.warning("No CPI-US All Items data available to display.")
    
        # Ensure PPI data exists and is formatted correctly
    if not df_ppi.empty:
            fig.add_trace(
                go.Scatter(
                    x=df_ppi['Month & Year'],
                    y=df_ppi['PPI Value'],
                    mode='lines',
                    name='PPI-US',
                    line=dict(color='#1C798A')
                )
            )
    else:
            st.warning("No PPI data available to display.")
    
        # Update Layout
    fig.update_layout(
            title='CPI & PPI Comparison',
            xaxis=dict(
                showgrid=False,
                showticklabels=True,
                color="#474747",
                tickfont=dict(color="#474747"),
            ),
            yaxis=dict(
                title='Index Value',
                showgrid=False,
                color="#474747",
                tickfont=dict(color="#474747"),
            ),
            legend=dict(
                orientation="h",
                x=0.01,
                y=-0.3,
                xanchor='left',
                yanchor='bottom',
                bgcolor='rgba(255, 255, 255, 0.6)',
                font=dict(size=10, color="#474747"),
            ),
            hovermode='x unified',
            plot_bgcolor='rgba(0,0,0,0)',
            paper_bgcolor='rgba(0,0,0,0)',
            height=300,
            width=600,
            margin=dict(b=60, t=20, r=15, l=15),
        )
    
    return fig

def plot_gdp_and_industry(selected_industry=None):
    fig = make_subplots(specs=[[{"secondary_y": True}]])

    # 1. Add GDP Value Line (Primary Axis)
    fig.add_trace(
        go.Scatter(
            x=df_gdp_filtered['Year'],
            y=df_gdp_filtered['Value'],
            mode='lines',
            name='GDP-Value',
            fill='tozeroy',  # Create area chart by filling to the x-axis
            fillcolor='#032649', #'rgba(235, 137, 40, 0.6)', 
            line=dict(color='#032649', width=2),
            marker=dict(size=10)
        ),
        secondary_y=False
    )

    # 2. Add GDP Percent Change Line (Secondary Axis)
    fig.add_trace(
        go.Scatter(
            x=df_gdp_filtered['Year'],
            y=df_gdp_filtered['Percent Change'],
            mode='lines',
            name='GDP-% Change',
            line=dict(color='#A5A5A5', width=2, dash='solid'),
            marker=dict(size=10)
        ),
        secondary_y=True
    )

    # Check if an industry is selected
    if selected_industry:
        df_industry_filtered = df_combined[df_combined['Industry'] == selected_industry]

        # 3. Add Selected Industry Value Line (Primary Axis)
        fig.add_trace(
            go.Scatter(
                x=df_industry_filtered['Year'],
                y=df_industry_filtered['Value'],
                mode='none',
                name=f'{selected_industry}-Value',
                fill='tozeroy',  # Area chart
                fillcolor='#EB8928', 
                line=dict(color='#EB8928', width=2),
                marker=dict(size=10)
            ),
            secondary_y=False
        )

        # 4. Add Selected Industry Percent Change Line (Secondary Axis)
        fig.add_trace(
            go.Scatter(
                x=df_industry_filtered['Year'],
                y=df_industry_filtered['Percent Change'],
                mode='lines',
                name=f'{selected_industry}-(% Change)',
                line=dict(color='#1C798A', width=2, dash='solid'),
                marker=dict(size=10)
            ),
            secondary_y=True
        )

    # Update layout
    fig.update_layout(
        title='',
        xaxis_title='',
        yaxis_title='Value',
        yaxis2_title='% Change',
        xaxis=dict(
            showgrid=False,
            showticklabels=True,
            color="#474747",  # X-axis label and line color
            tickfont=dict(color="#474747"),  # X-axis tick labels color
        ),
        yaxis=dict(
            title='',
            showgrid=False,
            color="#474747",  # Y-axis label and line color
            tickfont=dict(color="#474747"),  # Y-axis tick labels color
            tickformat=',',
        ),
        legend=dict(
            orientation="h",
            x=0.01,
            y=-0.15,
            xanchor='left', 
            yanchor='bottom',
            bgcolor='rgba(255, 255, 255, 0.6)',
            font=dict(size=10),
            traceorder='normal'
        ),
        template='plotly_white',
        plot_bgcolor='rgba(0,0,0,0)',
        paper_bgcolor='rgba(0,0,0,0)',
        height=450,
        width=700,
        margin=dict(b=60, t=20, r=15, l=15),  # Increased bottom margin for space
    )
    st.plotly_chart(fig, use_container_width=True)
    return fig

# Function to export charts to PowerPoint
def update_figure_slide(ppt, title, fig, slide_number, width, height, left, top):
    if fig is None:
        print(f"Skipping slide '{title}' because the figure is None.")
        return  # Skip if fig is None

    # Get the slide corresponding to the slide number (index starts at 0)
    slide = ppt.slides[slide_number]  # Adjust for 0-based index

    # # Set slide title (optionally adjust placement based on the layout)
    # title_shape = slide.shapes.title
    # title_shape.text = f"Slide {slide_number}: {title}"  # Add slide number to title

    # Save the figure image to a BytesIO object (no size, position parameters here)
    fig_image = BytesIO()
    fig.write_image(fig_image, format="png")  # Only pass the format here
    fig_image.seek(0)

    # Use Inches for size and position only in the add_picture() function
    slide.shapes.add_picture(fig_image, Inches(left), Inches(top), Inches(width), Inches(height))
    fig_image.close()

def export_all_to_pptx(labor_fig, external_fig, gdp_fig, cpi_ppi_fig):
    # Load the custom template
    template_path = os.path.join(os.getcwd(), "streamlit_dashboard", "data", "main_template_pitch.pptx")
    ppt = Presentation(template_path)  # Load the template

    # Use the existing slides (slide_number corresponds to the slide index)
    update_figure_slide(ppt, "labor Force & Unemployment Data", labor_fig, slide_number=5, width=5, height=2.50, left=0.08, top=1.3)
    update_figure_slide(ppt, "External Driver Indicators", external_fig, slide_number=7, width=4.50, height=3.75, left=5.20, top=1.3)
    update_figure_slide(ppt, "GDP by Industry", gdp_fig, slide_number=5, width=5.00, height=2.50, left=0.08, top=4.4)
    update_figure_slide(ppt, "CPI and PPI Comparison", cpi_ppi_fig, slide_number=5, width=4.55, height=2.50, left=5.10, top=1.3)

    # Save the PPT file to BytesIO and return the bytes
    ppt_bytes = BytesIO()
    ppt.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes

def get_us_indicators_layout():
    """Render the full dashboard layout and export data directly without session state."""
    st.set_page_config(page_title="US Indicators", layout="wide")
    # labor Force & Unemployment Data
    st.subheader("Labor Force Partication & Unemployment Rate with Population")
    labor_fig = plot_labor_unemployment()

    # External Driver Indicators
    st.subheader("External Driver Indicators")
    selected_indicators = st.multiselect(
        "Select External Indicators",
        options=external_driver_df["Indicator"].unique(),
        default=["World GDP"],
        key="external_indicators_multiselect"
    )
    external_fig = plot_external_driver(selected_indicators)

    # GDP by Industry
    st.subheader("GDP by Industry")
    selected_gdp_industry = st.selectbox(
        "Select Industry",
        options=gdp_industry,
        index=0,
        key="gdp_industry_selectbox"
    )
    gdp_fig = plot_gdp_and_industry(selected_gdp_industry)

    # CPI and PPI Comparison
    st.subheader("CPI & PPI")
    if cpi_industry:
        selected_cpi_industry = st.selectbox("Select CPI Industry", cpi_industry, index=0, key="cpi_industry_selectbox")
        cpi_ppi_fig = plot_cpi_ppi(selected_cpi_industry)
        st.plotly_chart(cpi_ppi_fig, use_container_width=True)
    else:
        st.warning("No industries available for selection.")

    if st.button("Export Charts to PowerPoint", key="export_button"):
        # Export the charts to PowerPoint using the export_all_to_pptx function
        pptx_file = export_all_to_pptx(labor_fig, external_fig, gdp_fig, cpi_ppi_fig)
        
        # Create a download button for the user to download the PowerPoint file
        st.download_button(
            label="Download PowerPoint",  # The label for the button
            data=pptx_file,  # The PowerPoint file content
            file_name="US_indicators.pptx",  # The default filename for the download
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"  # MIME type for PowerPoint
        )

get_us_indicators_layout()
