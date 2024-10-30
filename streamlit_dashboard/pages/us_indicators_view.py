import os
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from PIL import Image

# Define URLs and Paths
country = "USA"
timefrom = 2010
url_lfs = f"https://rplumber.ilo.org/data/indicator/?id=EAP_DWAP_SEX_AGE_RT_M&ref_area={country}&sex=SEX_T&classif1=AGE_AGGREGATE_TOTAL&timefrom={timefrom}&type=label&format=.csv"
url_unemp = f"https://rplumber.ilo.org/data/indicator/?id=UNE_DEAP_SEX_AGE_RT_M&ref_area={country}&sex=SEX_T&classif1=AGE_AGGREGATE_TOTAL&timefrom={timefrom}&type=label&format=.csv"
url_pop = "https://fred.stlouisfed.org/graph/fredgraph.csv?bgcolor=%23e1e9f0&chart_type=line&drp=0&fo=open%20sans&graph_bgcolor=%23ffffff&height=450&mode=fred&recession_bars=on&txtcolor=%23444444&ts=12&tts=12&width=1140&nt=0&thu=0&trc=0&show_legend=yes&show_axis_titles=yes&show_tooltip=yes&id=POPTHM&scale=left&cosd=2009-12-29&coed=2024-08-01&line_color=%234572a7&link_values=false&line_style=solid&mark_type=none&mw=3&lw=2&ost=-99999&oet=99999&mma=0&fml=a&fq=Monthly&fam=avg&fgst=lin&fgsnd=2020-02-01&line_index=1&transformation=lin&vintage_date=2024-10-09&revision_date=2024-10-09&nd=1959-01-01"
url_gdp_us = "https://apps.bea.gov/industry/Release/XLS/GDPxInd/GrossOutput.xlsx"
xls = pd.ExcelFile(url_gdp_us)

    # Labour Force Participation Rate Data
df_lfs = pd.read_csv(url_lfs)
df_lfs = df_lfs.rename(columns={'ref_area.label': 'country', 'obs_value': 'labour_force_rate'})
df_lfs['time'] = df_lfs['time'].astype(str)
time_split = df_lfs['time'].str.split('M', expand=True)
df_lfs['year'] = pd.to_numeric(time_split[0], errors='coerce').astype('Int64')
df_lfs['month'] = pd.to_numeric(time_split[1], errors='coerce').astype('Int64')

# Unemployment Rate Data
df_unemp = pd.read_csv(url_unemp)
df_unemp = df_unemp.rename(columns={'ref_area.label': 'country', 'obs_value': 'unemployment_rate'})
df_unemp['time'] = df_unemp['time'].astype(str)
time_split_unemp = df_unemp['time'].str.split('M', expand=True)
df_unemp['year'] = pd.to_numeric(time_split_unemp[0], errors='coerce').astype('Int64')
df_unemp['month'] = pd.to_numeric(time_split_unemp[1], errors='coerce').astype('Int64')

# Population Data
df_pop = pd.read_csv(url_pop)
df_pop = df_pop.rename(columns={'DATE': 'date', 'POPTHM': 'population'})
df_pop['date'] = pd.to_datetime(df_pop['date'])
df_pop['year'] = df_pop['date'].dt.year
df_pop['month'] = df_pop['date'].dt.month

# External Driver Data
external_driver_path = r"streamlit_dashboard/data/business_enviornmental_profiles_final.xlsx"
external_driver_df = pd.read_excel(external_driver_path)
external_driver_df['Year'] = pd.to_numeric(external_driver_df['Year'], errors='coerce')
indicator_mapping = {indicator: {'label': indicator, 'value': indicator} for indicator in external_driver_df['Indicator'].unique()}
external_driver_df['Indicator_Options'] = external_driver_df['Indicator'].map(indicator_mapping)

    # CPI Industry Mapping
industry_mapping = {
    'All items': 'CUSR0000SA0',
    'Energy': 'CUSR0000SA0E',
    'Apparel': 'CUSR0000SAA',
    'Medical care': 'CUSR0000SAM',
    'Housing': 'CUSR0000SAH',
    'Food': 'CUSR0000SAF1'
    }

file_path = r"streamlit_dashboard/data/CPI_industry.txt"
ppi_file_path = r"streamlit_dashboard/data/PPI.txt"
# Load CPI data
df = pd.read_csv(file_path, delimiter=',').dropna().reset_index(drop=True)
df_unpivoted = df.melt(id_vars=["Series ID"], var_name="Month & Year", value_name="Value")
df_unpivoted = df_unpivoted[df_unpivoted["Value"].str.strip() != ""]
df_unpivoted["Series ID"] = df_unpivoted["Series ID"].astype(str)
df_unpivoted["Value"] = pd.to_numeric(df_unpivoted["Value"], errors='coerce')
df_unpivoted["Month & Year"] = pd.to_datetime(df_unpivoted["Month & Year"], format='%b %Y', errors='coerce')
df_cleaned = df_unpivoted.dropna(subset=["Series ID", "Month & Year", "Value"])
all_items_data = df_cleaned[df_cleaned['Series ID'] == 'CUSR0000SA0']
all_items_data = all_items_data[all_items_data['Month & Year'] >= '2010-01-01']
    # Function to fetch CPI data for the selected industry

    # Load and clean PPI data
df_ppi = pd.read_csv(ppi_file_path, delimiter=',').dropna().reset_index(drop=True)
df_ppi_unpivoted = df_ppi.melt(id_vars=["Year"], var_name="Month", value_name="Value")
df_ppi_unpivoted["Month & Year"] = pd.to_datetime(df_ppi_unpivoted["Month"] + " " + df_ppi_unpivoted["Year"].astype(str),format='%b %Y', errors='coerce')
df_ppi_unpivoted['Value'] = pd.to_numeric(df_ppi_unpivoted['Value'], errors='coerce')
df_ppi_unpivoted = df_ppi_unpivoted.dropna(subset=['Month & Year', 'Value'])
df_ppi_unpivoted = df_ppi_unpivoted[df_ppi_unpivoted["Month & Year"] >= '2010-01-01']

    # Clean and reshape GDP data
df_gdp_us = pd.read_excel(xls, sheet_name="TGO105-A")
df_gdp_us = df_gdp_us.iloc[6:].reset_index(drop=True)
df_gdp_us.columns = df_gdp_us.iloc[0]
df_gdp_us = df_gdp_us.drop(0).reset_index(drop=True)
df_gdp_us = df_gdp_us.drop(columns=["Line"])
df_gdp_us = df_gdp_us.drop(df_gdp_us.columns[1], axis=1)
df_gdp_us = df_gdp_us.rename(columns={df_gdp_us.columns[df_gdp_us.isna().any()].tolist()[0]: 'Industry'})
df_gdp_us["Industry"] = df_gdp_us["Industry"].replace("    All industries", "GDP")
df_gdp_us["Industry"] = df_gdp_us["Industry"].str.replace("  ", "")
df_gdp_unpivoted = df_gdp_us.melt(id_vars=["Industry"], var_name="Year", value_name="Value")
df_gdp_unpivoted["Year"] = df_gdp_unpivoted["Year"].astype(int)
df_gdp_unpivoted["Value"] = pd.to_numeric(df_gdp_unpivoted["Value"], errors='coerce')
df_gdp_unpivoted = df_gdp_unpivoted.dropna(subset=["Value"])

    # Clean and reshape GDP Percent Change data
df_pct_change = pd.read_excel(xls, sheet_name="TGO101-A")
df_pct_change = df_pct_change.iloc[6:].reset_index(drop=True)
df_pct_change.columns = df_pct_change.iloc[0]
df_pct_change = df_pct_change.drop(0).reset_index(drop=True)
df_pct_change = df_pct_change.drop(columns=["Line"])
df_pct_change = df_pct_change.drop(df_pct_change.columns[1], axis=1)
df_pct_change = df_pct_change.rename(columns={df_pct_change.columns[df_pct_change.isna().any()].tolist()[0]: 'Industry'})
df_pct_change["Industry"] = df_pct_change["Industry"].replace("    All industries", "GDP")
df_pct_change["Industry"] = df_pct_change["Industry"].str.replace("  ", "")
df_pct_unpivoted = df_pct_change.melt(id_vars=["Industry"], var_name="Year", value_name="Percent Change")
df_pct_unpivoted["Year"] = df_pct_unpivoted["Year"].astype(int)
df_pct_unpivoted["Percent Change"] = pd.to_numeric(df_pct_unpivoted["Percent Change"], errors='coerce')
df_pct_unpivoted = df_pct_unpivoted.dropna(subset=["Percent Change"])

df_combined = pd.merge(
    df_gdp_unpivoted,
    df_pct_unpivoted,
    on=["Industry", "Year"],
    how="inner"
   )

    # Filter GDP data
df_gdp_filtered = df_combined[df_combined['Industry'] == 'GDP']

    # Create a list of industries excluding GDP for the dropdown
industry_options = df_combined['Industry'].unique().tolist()
industry_options.remove('GDP')

def fetch_cpi_data(series_id, df_cleaned):
    selected_data = df_cleaned[df_cleaned['Series ID'] == series_id]
    selected_data = selected_data[selected_data['Month & Year'] >= '2010-01-01']
    return selected_data[['Month & Year', 'Value']].rename(columns={'Month & Year': 'date', 'Value': 'value'})

# if 'counter' not in st.session_state:
#     st.session_state.counter = 0

def plot_labour_unemployment():
    # Merge unemployment and labour force data
    merged = pd.merge(df_lfs, df_unemp, on=["year", "month", "country"], how='inner')
    merged = pd.merge(merged, df_pop, on=["year", "month"], how='inner')

    fig = go.Figure()

    # Plot population as an area chart on the primary y-axis
    min_population = merged['population'].min()
    fig.add_trace(go.Scatter(
        x=pd.to_datetime(merged[['year', 'month']].assign(day=1)),
        y=merged['population'],
        fill='tozeroy',  # Area chart
        name='Population',
        mode='lines',
        line=dict(color='blue'),
        yaxis='y1'
    ))

    # Plot unemployment rate on the secondary y-axis
    fig.add_trace(go.Scatter(
        x=pd.to_datetime(merged[['year', 'month']].assign(day=1)),
        y=merged['unemployment_rate'],
        name='Unemployment Rate',
        mode='lines',
        line=dict(color='red'),
        yaxis='y2'
    ))

    # Plot labour force participation rate on the secondary y-axis
    fig.add_trace(go.Scatter(
        x=pd.to_datetime(merged[['year', 'month']].assign(day=1)),
        y=merged['labour_force_rate'],
        name='Labour Force Participation Rate',
        mode='lines',
        line=dict(color='green'),
        yaxis='y2'
    ))

    fig.update_layout(
        title='Population, Unemployment Rate, and Labour Force Participation Rate (USA)',
        xaxis=dict(title='Date'),
        yaxis=dict(
            title='Population',
            side='left',
            range=[min_population, merged['population'].max() * 1.1]
        ),
        yaxis2=dict(
            title='Rate (%)',
            overlaying='y', 
            side='right'
        ),
        legend=dict(
            x=0.01, y=0.99, bgcolor='rgba(255, 255, 255, 0.6)',
            bordercolor='black', borderwidth=1
        ),
        hovermode='x'
    )

    # key = f"labour_unemployment_chart_{st.session_state.counter}"
    # st.session_state.counter += 1 

    st.plotly_chart(fig, use_container_width=True)


def plot_external_driver(selected_indicators):

    if not selected_indicators:
        selected_indicators = ["World GDP"]

    fig = go.Figure()

    for indicator in selected_indicators:
        indicator_data = external_driver_df[external_driver_df['Indicator'] == indicator]

        if '% Change' not in indicator_data.columns:
            raise ValueError(f"Expected '% Change' column not found in {indicator}")

        fig.add_trace(
            go.Scatter(
                x=indicator_data['Year'],
                y=indicator_data['% Change'],
                mode='lines',
                name=indicator
            )
        )

    fig.update_layout(
        title='External Driver Indicators',
        xaxis=dict(title=''),
        yaxis=dict(title='Percent Change'),
        hovermode='x'
    )
    # key = f"external_driver_chart_{st.session_state.counter}"
    # st.session_state.counter += 1

    st.plotly_chart(fig, use_container_width=True)#, key=key)


def plot_cpi_ppi(selected_series_id):
    """
    Plot CPI and PPI data on a single chart for comparison.
    """
    fig = go.Figure()

    # Fetch and plot the selected CPI industry data
    cpi_data = fetch_cpi_data(selected_series_id, df_cleaned)
    if not cpi_data.empty:
        fig.add_trace(
            go.Scatter(
                x=cpi_data['date'],
                y=cpi_data['value'],
                mode='lines',
                name='CPI by Industry',
                line=dict(color='blue')
            )
        )
    else:
        st.warning(f"No data available for the selected CPI series: {selected_series_id}")

    # Plot CPI-US All Items data
    if not all_items_data.empty:
        fig.add_trace(
            go.Scatter(
                x=all_items_data['Month & Year'],
                y=all_items_data['Value'],
                mode='lines',
                name='CPI-US (All Items)',
                line=dict(color='green', dash='dash')
            )
        )
    else:
        st.warning("No CPI-US All Items data available to display.")

    # Plot aggregated PPI data
    if not df_ppi_unpivoted.empty:
        df_ppi_aggregated = df_ppi_unpivoted.groupby('Month & Year', as_index=False).agg({'Value': 'mean'})
        fig.add_trace(
            go.Scatter(
                x=df_ppi_aggregated['Month & Year'],
                y=df_ppi_aggregated['Value'],
                mode='lines',
                name='PPI-US',
                line=dict(color='red')
            )
        )
    else:
        st.warning("No PPI data available to display.")

    # Configure the layout of the chart
    fig.update_layout(
        title='CPI and PPI Comparison',
        xaxis=dict(title='Date'),
        yaxis=dict(title='Value'),
        hovermode='x unified'
    )

    # key = f"cpi_ppi_chart_{st.session_state.counter}"
    # st.session_state.counter += 1 

    st.plotly_chart(fig, use_container_width=True)#, key=key)


def plot_gdp_and_industry(selected_industry=None):
    fig = make_subplots(specs=[[{"secondary_y": True}]])

    # 1. Add GDP Value Line (Primary Axis)
    fig.add_trace(
        go.Scatter(
            x=df_gdp_filtered['Year'],
            y=df_gdp_filtered['Value'],
            mode='lines',
            name='GDP - Value',
            line=dict(color='blue', width=2, dash='solid'),
            marker=dict(size=6)
        ),
        secondary_y=False
    )

    # 2. Add GDP Percent Change Line (Secondary Axis)
    fig.add_trace(
        go.Scatter(
            x=df_gdp_filtered['Year'],
            y=df_gdp_filtered['Percent Change'],
            mode='lines',
            name='GDP - Percent Change',
            line=dict(color='orange', width=2, dash='solid'),
            marker=dict(size=6)
        ),
        secondary_y=True
    )

    # Check if an industry is selected
    if selected_industry:
        df_industry_filtered = df_combined[df_combined['Industry'] == selected_industry]

        # 3. Add Selected Industry Value Line (Primary Axis)
        fig.add_trace(
            go.Scatter(
                x=df_industry_filtered['Year'],
                y=df_industry_filtered['Value'],
                mode='lines',
                name=f'{selected_industry} - Value',
                line=dict(color='red', width=2, dash='solid'),
                marker=dict(size=6)
            ),
            secondary_y=False
        )

        # 4. Add Selected Industry Percent Change Line (Secondary Axis)
        fig.add_trace(
            go.Scatter(
                x=df_industry_filtered['Year'],
                y=df_industry_filtered['Percent Change'],
                mode='lines',
                name=f'{selected_industry} - Percent Change',
                line=dict(color='green', width=2, dash='solid'),
                marker=dict(size=6)
            ),
            secondary_y=True
        )

    # Update layout
    fig.update_layout(
        title=f'GDP and {selected_industry or "GDP"} - Value & Percent Change over Years',
        xaxis_title='',
        yaxis_title='Value',
        yaxis2_title='Percent Change',
        legend_title='Legend',
        template='plotly_white'
    )

    # key = f"gdp_chart_{st.session_state.counter}"
    # st.session_state.counter += 1 

    st.plotly_chart(fig, use_container_width=True)#, key=key)
def export_all_to_pptx(labour_image, external_driver_image, gdp_image, cpi_ppi_image):
    pptx_io = BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    for img in [labour_image, external_driver_image, gdp_image, cpi_ppi_image]:
        img_buf = BytesIO()
        img.save(img_buf, format='PNG')
        img_buf.seek(0)
        slide.shapes.add_picture(img_buf, Inches(1), Inches(1), width=Inches(6))

    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def get_us_indicators_layout():
    """Render the full dashboard layout."""
    st.title("US Indicators Dashboard")

    # Labour Force and Unemployment Section
    st.subheader("Labour Force & Unemployment Data")
    labour_fig = plot_labour_unemployment()

    # External Driver Indicators Section
    st.subheader("External Driver Indicators")
    selected_indicators = st.multiselect(
        "Select External Indicators",
        options=external_driver_df["Indicator"].unique(),
        default=["World GDP"],
        key="external_indicators_multiselect"
    )
    external_fig = plot_external_driver(selected_indicators)

    # GDP by Industry Section
    st.subheader("GDP by Industry")
    selected_gdp_industry = st.selectbox(
        "Select Industry",
        options=df_combined["Industry"].unique(),
        index=0,
        key="gdp_industry_selectbox"
    )
    gdp_fig = plot_gdp_and_industry(selected_gdp_industry)

    # CPI and PPI Comparison Section
    st.subheader("CPI and PPI Comparison")
    selected_cpi_series = st.selectbox(
        "Select CPI Industry", 
        list(industry_mapping.keys()), 
        index=1, 
        key="cpi_series_selectbox"
    )
    selected_series_id = industry_mapping[selected_cpi_series]
    cpi_ppi_fig = plot_cpi_ppi(industry_mapping[selected_cpi_series])

    if st.button("Export All Charts to PPTX", key="export_button"):

        dummy_image = Image.new("RGB", (400, 300), color="white")
        pptx_file = export_all_to_pptx(
            dummy_image, dummy_image, dummy_image, dummy_image
        )
        st.download_button(
            label="Download PPTX",
            data=pptx_file,
            file_name="us_indicators_dashboard.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
