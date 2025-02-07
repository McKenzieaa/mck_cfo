import streamlit as st
import pandas as pd
import requests
import plotly.express as px
import zipfile
import io
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import plotly.io as pio
from io import StringIO
import os
import mysql.connector
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import boto3
from botocore.exceptions import NoCredentialsError, ClientError

st.set_page_config(page_title="Global Industry Analysis", layout="wide")

# MySQL database connection details 
host = st.secrets["mysql"]["host"]
user = st.secrets["mysql"]["user"]
password = st.secrets["mysql"]["password"]
database = st.secrets["mysql"]["database"]

# AWS S3 bucket details
bucket_name = 'documentsapi'
energy_prices_df = 'industry_data/energy_data/2._Energy_Prices.csv'
sales_price_df = 'industry_data/energy_data/7a._U.S._Electricity_Industry_Overview.csv'

def read_csv_from_s3(bucket, key):
    try:
        # Check if all AWS secrets are present
        if not all(key in st.secrets["aws"] for key in ["AWS_ACCESS_KEY_ID", "AWS_SECRET_ACCESS_KEY", "AWS_DEFAULT_REGION"]):
            st.error("AWS credentials are missing or incomplete in secrets.toml.")
            return None

        s3 = boto3.client(
            's3',
            aws_access_key_id=st.secrets["aws"]["AWS_ACCESS_KEY_ID"],
            aws_secret_access_key=st.secrets["aws"]["AWS_SECRET_ACCESS_KEY"],
            region_name=st.secrets["aws"]["AWS_DEFAULT_REGION"]
        )
        response = s3.get_object(Bucket=bucket, Key=key)
        content = response['Body'].read().decode('utf-8')
        return StringIO(content)
    except KeyError as e:
        st.error(f"Missing key in secrets.toml: {e}")
        return None
    except NoCredentialsError:
        st.error("AWS credentials not found. Please configure your credentials.")
        return None
    except Exception as e:
        st.error(f"Error reading from S3: {e}")
        return None

# Read the CSV files from S3
csv_file1 = read_csv_from_s3(bucket_name, energy_prices_df)
csv_file2 = read_csv_from_s3(bucket_name, sales_price_df)

if csv_file1 is None or csv_file2 is None:
    st.stop()  # Stop execution if CSVs couldn't be loaded

# Process the first CSV (Energy Prices)
try:
    df1 = pd.read_csv(csv_file1, skiprows=4)
    columns_to_drop = [
        'map', 'linechart', 'source key', 'Unnamed: 1', '1960', '1961', '1962', 
        '1963', '1964', '1965', '1966', '1967', '1968', '1969', '1970', '1971', 
        '1972', '1973', '1974', '1975', '1976', '1977', '1978', '1979', '1980', 
        '1981', '1982', '1983', '1984', '1985', '1986', '1987', '1988', '1989'
    ]
    df1 = df1.drop(columns=[col for col in columns_to_drop if col in df1.columns])
    df1 = df1.rename(columns={'remove': 'Prices to Ultimate Customers', 'units': 'Units'})
    df1['Prices to Ultimate Customers'] = df1['Prices to Ultimate Customers'].str.replace(' Sector', '', regex=False)
    price_customers = df1.tail(3)
    id_vars = ['Prices to Ultimate Customers', 'Units']
    price_customers = pd.melt(price_customers, id_vars=id_vars, var_name='Year', value_name='Value')
except Exception as e:
    st.error(f"Error processing energy prices data: {e}")
    st.stop()

# Process the second CSV (Sales Prices)
try:
    df2 = pd.read_csv(csv_file2, skiprows=4)
    df2 = df2.drop(columns=[col for col in columns_to_drop if col in df2.columns])
    df2 = df2.rename(columns={'remove': 'Sales to Ultimate Customers', 'units': 'Units'})
    df2['Sales to Ultimate Customers'] = df2['Sales to Ultimate Customers'].str.replace(' Sector', '', regex=False)
    sales_customers = df2.iloc[16:-27].reset_index(drop=True)
    id_vars = ['Sales to Ultimate Customers', 'Units']
    sales_customers1 = pd.melt(sales_customers, id_vars=id_vars, var_name='Year', value_name='Value')
except Exception as e:
    st.error(f"Error processing sales prices data: {e}")
    st.stop()

# Connect to the MySQL database
try:
    conn = mysql.connector.connect(
        host=host,
        user=user,
        password=password,
        database=database
    )
except mysql.connector.Error as e:
    st.error(f"Error connecting to MySQL: {e}")
    st.stop()

query = """
SELECT 
    `Year`, `Target`, `EV/Revenue`, `EV/EBITDA`, `Business Description`, `Industry`, `Location`
FROM 
    precedent_table
"""

try:
    df_pt = pd.read_sql(query, conn)
except Exception as e:
    st.error(f"Error loading data from MySQL: {e}")
    st.stop()


query3 = """
SELECT `NAICS`, `LineItems`, `Percent`, `ReportID`, `Industry`, `Value`
FROM rma_table
"""
try:
    df_rma = pd.read_sql(query3, conn)
except Exception as e:
    st.error(f"Error loading data from MySQL (Public Companies): {e}")
    st.stop()

conn.close()

# Precdent Transaction
industries_pt = [
    'Hydroelectric Power Generation',
    'Natural Gas Extraction',
    'Petroleum and Petroleum Products Merchant Wholesalers (except Bulk Stations and Terminals)',
    'Crude Petroleum Extraction',
    # 'Pipeline Transportation of Natural Gas',
    'Natural Gas Distribution',
    # 'Industrial Gas Manufacturing'
]
df_pt_filter = df_pt[df_pt['Industry'].isin(industries_pt)]
df_pt_grouped = df_pt_filter.groupby('Year')[['EV/Revenue', 'EV/EBITDA']].mean().reset_index()
df_pt_grouped['EV/Revenue'] = df_pt_grouped['EV/Revenue'].round(2)
df_pt_grouped['EV/EBITDA'] = df_pt_grouped['EV/EBITDA'].round(2)

# RMA Data
df_rma_filtered = df_rma[df_rma['NAICS'].astype(str).str.startswith('2211')]
df_rma_filtered = df_rma_filtered.groupby(['ReportID', 'LineItems'], as_index=False)['Value'].mean()
df_rma_filtered['Value'] = (df_rma_filtered['Value'] / 1_000_000).round(2)
df_rma_is = df_rma_filtered[df_rma_filtered['ReportID'] == 'Income Statement']
df_rma_bs = df_rma_filtered[df_rma_filtered['ReportID'].isin(['Assets', 'Liabilities & Equity'])]
df_rma_bs_grouped = (df_rma_bs.groupby(['ReportID', 'LineItems'], as_index=False)['Value'].mean())
df_rma_bs_grouped['Grouped_LineItems'] =  df_rma_bs_grouped['LineItems']


# Selected Countries
selected_countries = ['China', 'India', 'World', 'Japan','Brazil','France', 'United States']

# Chart 1 = Market Size
category_data = [
    ('22T', 'Utilities', 'QREV', 'QSS'),
    ('2211T', 'Electric Power Generation, Transmission and Distribution', 'QREV', 'QSS')
]

categories = {
    category.strip(): {"cat_code": cat_code, "data_type": data_type, "program": program}
    for cat_code, category, data_type, program in category_data
}

selected_categories = ['Electric Power Generation, Transmission and Distribution']

def get_market_size_data(selected_categories):
    """Fetch and prepare market size data based on selected categories."""
    all_data = pd.DataFrame()

    for category in selected_categories:
        category = category.strip()
        try:
            cat_data = categories[category]
        except KeyError:
            continue

        url = f"https://www.census.gov/econ_export/?format=xls&mode=report&default=false&errormode=Dep&charttype=&chartmode=&chartadjn=&submit=GET+DATA&program={cat_data['program']}&startYear=2014&endYear=2024&categories%5B0%5D={cat_data['cat_code']}&dataType={cat_data['data_type']}&geoLevel=US&adjusted=false&notAdjusted=true&errorData=false&vert=1"
        response = requests.get(url)
        if response.status_code != 200:
            continue

        marketsize_data = BytesIO(response.content)
        df = pd.read_excel(marketsize_data)
        df = df.iloc[6:].reset_index(drop=True)
        df.columns = df.iloc[0]
        df = df.drop(0).reset_index(drop=True)
        df = df.dropna(axis=1, how='all')
        df['Value'] = pd.to_numeric(df['Value'].replace("N", pd.NA), errors='coerce')
        df[['Quarter', 'Year']] = df['Period'].str.split('-', expand=True)
        df['Quarter'] = df['Quarter'].replace({
            '1st Quarter': 'Q1', '2nd Quarter': 'Q2',
            '3rd Quarter': 'Q3', '4th Quarter': 'Q4'
        })
        df['Quarter_Year'] = df['Quarter'] + ' ' + df['Year'].astype(str)
        df['Category'] = category
        all_data = pd.concat([all_data, df], ignore_index=True)

    if all_data.empty:
        st.write("No data available for the selected categories.")
        return None

    return all_data

# Fetch Energy Data
url = "https://www.eia.gov/totalenergy/data/monthly/Zip_Excel_Month_end/MER_2024_09.zip"
response = requests.get(url)
url2 = "https://nyc3.digitaloceanspaces.com/owid-public/data/energy/owid-energy-data.csv"

with zipfile.ZipFile(io.BytesIO(response.content)) as z:
    with z.open("Table 07.01.xlsx") as f1:
        df_electricity_end_use = pd.read_excel(f1, sheet_name="Annual Data", skiprows=8)
    with z.open("Table 09.08.xlsx") as f2:
        df_avg_price = pd.read_excel(f2, sheet_name="Annual Data", skiprows=8)

# Chart 2 = electricity_end_use & avg_price (chart pending)
df_electricity_end_use.rename(columns={df_electricity_end_use.columns[0]: "Year"}, inplace=True)
df_avg_price.rename(columns={df_avg_price.columns[0]: "Year"}, inplace=True)

df_electricity_gen = pd.read_csv(url2)
df_renew_share = df_electricity_gen.dropna(subset=['renewables_share_elec'])


df_per_cap_elec_gen = df_electricity_gen.dropna(subset=['fossil_elec_per_capita', 'nuclear_elec_per_capita', 'renewables_elec_per_capita'])
df_per_cap_elec_gen = df_per_cap_elec_gen[df_per_cap_elec_gen['year'] == 2023]

selected_countries = ['China', 'India', 'World', 'Japan','Brazil','France', 'United States']
df_per_cap_elec_gen = df_per_cap_elec_gen[df_per_cap_elec_gen['country'].isin(selected_countries)]

df_per_cap_elec_gen = df_per_cap_elec_gen.melt(
    id_vars=['country'],
    value_vars=['fossil_elec_per_capita', 'nuclear_elec_per_capita', 'renewables_elec_per_capita'],
    var_name='Energy Source',
    value_name='Per Capita Generation'
)

df_per_cap_elec_gen['Energy Source'] = df_per_cap_elec_gen['Energy Source'].replace({
    'fossil_elec_per_capita': 'Fossil',
    'nuclear_elec_per_capita': 'Nuclear',
    'renewables_elec_per_capita': 'Renewables'
})
 
ele_gen_url = "https://www.eia.gov/totalenergy/data/browser/csv.php?tbl=T07.02B"
df_electricity_gen2 = pd.read_csv(ele_gen_url)
df_electricity_gen2['Description'] = df_electricity_gen2['Description'].str.extract(r'From (.*?),')
df_electricity_gen2 = df_electricity_gen2[df_electricity_gen2['Description'].notna()]
df_electricity_gen2.drop(columns=['MSN'], inplace=True)
df_electricity_gen2.drop(columns=['Column_Order'], inplace=True, errors='ignore')
df_electricity_gen2['Year'] = df_electricity_gen2['YYYYMM'].astype(str).str[:4]
df_electricity_gen2['YYYYMM'] = df_electricity_gen2['YYYYMM'].astype(str)
df_electricity_gen2 = df_electricity_gen2[df_electricity_gen2['YYYYMM'].str.endswith('13')]
df_electricity_gen2.drop(columns=['YYYYMM'], inplace=True)
source_category_mapping = {
    'Coal': 'Fossil Fuel',
    'Petroleum': 'Fossil Fuel',
    'Natural Gas': 'Fossil Fuel',
    'Other Fossil Gases': 'Fossil Fuel',
    'Nuclear Electric Power': 'Nuclear',
    'Hydroelectric Pumped Storage': 'Hydroelectric',
    'Conventional Hydroelectric Power': 'Hydroelectric',
    'Wood': 'Other',
    'Waste': 'Other',
    'Geothermal': 'Other',
    'Solar': 'Solar',
    'Wind': 'Wind'
}
df_electricity_gen2['Category'] = df_electricity_gen2['Description'].map(source_category_mapping)
# df_electricity_gen2 = df_electricity_gen2[df_electricity_gen2['Year']]
df_electricity_gen2['Value'] = pd.to_numeric(df_electricity_gen2['Value'], errors='coerce').mean()

# Energy Consumption
ene_cons = "https://www.eia.gov/totalenergy/data/browser/csv.php?tbl=T07.06"
df_ene_cons = pd.read_csv(ene_cons)
df_ene_cons['Description'] = df_ene_cons['Description'].str.split(',', n=1).str[1]
df_ene_cons = df_ene_cons[df_ene_cons['Description'].str.contains("Residential|Transportation|Industrial|Commercial", case=False, na=False)]
df_ene_cons['Year'] = df_ene_cons['YYYYMM'].astype(str).str[:4]
df_ene_cons = df_ene_cons[['Year', 'Description', 'Value']]
df_ene_cons['Value'] = df_ene_cons['Value'].round(1)
# df_ene_cons = df_ene_cons.groupby(['Year', 'Description'], as_index=False).sum()

# Share of electricity production from renewables
share_elec_prod = "https://ourworldindata.org/grapher/share-electricity-renewables.csv?v=1&csvType=full&useColumnShortNames=true"
# Fetch the data.
try:
    response = requests.get(share_elec_prod)
    response.raise_for_status() 
    csv_data = StringIO(response.text)
    df_share_elec_prod = pd.read_csv(csv_data)
except requests.exceptions.RequestException as e:
    print(f"Failed to fetch data: {e}")
    exit()

df_share_elec_prod.rename(columns={'Entity': 'Countries'}, inplace=True)
filt_share_elec_prod = df_share_elec_prod[df_share_elec_prod['Countries'].isin(selected_countries)]

if 'Year' in filt_share_elec_prod.columns and 'renewable_share_of_electricity__pct' in filt_share_elec_prod.columns:

    filt_share_elec_prod = filt_share_elec_prod.dropna(subset=['Year', 'renewable_share_of_electricity__pct'])
    filt_share_elec_prod['Year'] = pd.to_numeric(filt_share_elec_prod['Year'], errors='coerce')
    filt_share_elec_prod['renewable_share_of_electricity__pct'] = pd.to_numeric(
        filt_share_elec_prod['renewable_share_of_electricity__pct'], errors='coerce'
    )
    filt_share_elec_prod = filt_share_elec_prod.dropna(subset=['Year', 'renewable_share_of_electricity__pct'])
else:
    raise ValueError("Required columns 'Year' or 'renewable_share_of_electricity__pct' are missing from the DataFrame.")

per_cap_electricity = pd.read_csv("https://ourworldindata.org/grapher/per-capita-electricity-fossil-nuclear-renewables.csv?v=1&csvType=full&useColumnShortNames=true", 
                                  storage_options={'User-Agent': 'Our World In Data data fetch/1.0'})
per_cap_electricity = per_cap_electricity.rename(columns={
    'Entity': 'country',     
    per_cap_electricity.columns[3]: 'fossil',  
    per_cap_electricity.columns[4]: 'nuclear', 
    per_cap_electricity.columns[5]: 'renewable'
})

per_cap_electricity = per_cap_electricity.drop(columns=['Code'])
filter_per_cap_electricity = per_cap_electricity[(per_cap_electricity['Year'] == 2023) & 
                                    (per_cap_electricity['country'].isin(selected_countries))]

filter_per_cap_electricity[['fossil', 'nuclear', 'renewable']] = filter_per_cap_electricity[['fossil', 'nuclear', 'renewable']].div(
    filter_per_cap_electricity[['fossil', 'nuclear', 'renewable']].sum(axis=1), axis=0)


# ENERGY
st.markdown("<h2 style='font-weight: bold; font-size:24px;'>Energy</h2>", unsafe_allow_html=True)
with st.expander("", expanded=True): 
    market_data = get_market_size_data(selected_categories)
    if market_data is not None:
        yearly_data = market_data.groupby(['Year', 'Category'], as_index=False).agg({'Value': 'mean'})

        PRIMARY_COLORS = {
            'dark_blue': '#032649',
            'orange': '#EB8928',
            'dark_grey': '#595959',
            'light_grey': '#A5A5A5',
            'turquoise_blue': '#1C798A'
        }

        fig1 = px.bar(
            yearly_data,
            x='Year', y='Value', color='Category',
            title="Market Size",
            labels={'Value': 'Market-Size (in millions)', 'Year': ''},
            color_discrete_sequence=[PRIMARY_COLORS["dark_blue"]]
        )

        fig1.update_layout(
            legend=dict(
                x=0,  # Position at the left
                y=1,  # Position at the top
                title="",
                xanchor='left', 
                yanchor='top',
                font=dict(size=8)  # Set font size to 8
            ),
            plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',margin=dict(l=0, r=0, t=0) #,width=610,height=240
        )
        # st.plotly_chart(fig1)

        fig2 = make_subplots(specs=[[{"secondary_y": True}]])

        # Add bar chart for Electricity End Use to the primary axis
        fig2.add_trace(
            go.Bar(
                x=df_electricity_end_use["Year"],
                y=df_electricity_end_use[df_electricity_end_use.columns[1]],
                name="Electricity End Use (Billion Kilowatthours)",
                marker_color="#032649"
            ),
            secondary_y=False,
        )

        # Add line chart for Average Price of Electricity to the secondary axis
        fig2.add_trace(
            go.Scatter(
                x=df_avg_price["Year"],
                y=df_avg_price[df_avg_price.columns[1]],
                name="Average Price of Electricity (Cents per Kilowatthour)",
                mode="lines",
                line=dict(color="#FF5733")
            ),
            secondary_y=True,
        )

        # Update layout for titles and axes
        fig2.update_layout(
            title_text="Electricity End Use and Average Price of Electricity",
            xaxis_title="Year",
            yaxis_title="Electricity End Use (Billion Kilowatthours)",
            legend=dict(x=0.01, y=0.99),
        )

        # Set secondary y-axis title
        fig2.update_yaxes(
            title_text="Average Price of Electricity (Cents per Kilowatthour)", 
            secondary_y=True
        )
        # st.plotly_chart(fig3)

        # Electricity Generation Map
        # st.sidebar.header("Electricity Generation")
        selected_year =(2023) #st.sidebar.slider("Select Year", 2000, 2023, 2023)
        df_selected_year = df_electricity_gen[df_electricity_gen["year"] == selected_year]
        fig4 = px.choropleth(
            df_selected_year,
            locations='country',
            locationmode='country names',
            color='electricity_generation',
            title=f'Electricity Generation by Country ({selected_year})',
            labels={'electricity_generation': 'Electricity Generation (GWh)'},
            color_continuous_scale="Blues"  # Example color scale; you can choose others like "Plasma", "Blues", etc.
        )

        # st.plotly_chart(fig4)

        # Renewable Share of Electricity
        if selected_countries:
            filtered_df = df_renew_share[df_renew_share["country"].isin(selected_countries)]

        fig5 = px.line(
            filtered_df,
            x="year", y="renewables_share_elec", color="country",
            title="Renewable Share of Electricity"
        )
        colors_list = list(PRIMARY_COLORS.values())
        for i, trace in enumerate(fig5.data):
            trace.update(line=dict(color=colors_list[i % len(colors_list)])) 
        fig5.update_layout(
            plot_bgcolor='rgba(0,0,0,0)',
            paper_bgcolor='rgba(0,0,0,0)',
            margin=dict(l=10, r=10, t=10, b=10),
            # width=625,
            # height=250
        )
            # st.plotly_chart(fig5)

        # Solar Projects Coming Up Next 12 Months
        # st.sidebar.header("Map of Solar Projects Coming Up Next 12 Months")
        solar_url = "https://www.eia.gov/electricity/monthly/epm_table_grapher.php?t=table_6_05"
        solar_data = pd.read_html(solar_url)[1]
        st.dataframe(solar_data)

        # st.sidebar.header("Per Capita Electricity")
        fig6 = px.bar(
            df_per_cap_elec_gen,
            x='Per Capita Generation',
            y='country',
            color='Energy Source',
            orientation='h',  # Horizontal orientation
            title='Electricity Generation per Capita by Energy Source (Major Countries in 2023)',
            labels={'country': 'Country', 'Per Capita Generation': 'Percentage of Total Generation'},
            color_discrete_map={
                'Fossil': '#032649',        
                'Nuclear': '#1C798A',      
                'Renewables': '#EB8928'     
            }
        )
        fig6.update_layout(barmode='stack')
        fig6.update_xaxes(title_text="")
        # st.plotly_chart(fig6)

        df_ene_cons_2023 = df_ene_cons[df_ene_cons['Year'] == '2023']
        fig7 = px.pie(df_ene_cons_2023, 
                names='Description', 
                values='Value', 
                title='Energy Consumption by Source in 2023',
                labels={'Value': 'Energy Value', 'Description': 'Energy Source'})
    
        fig8 = px.bar(
            df_ene_cons,
            x='Year', y='Value', color='Description',
            title='Energy Source Distribution Over Years',
            labels={'Value': 'Energy Value', 'Year': 'Year'},
            color_discrete_map=PRIMARY_COLORS
        )


        fig10 = px.bar(df_electricity_gen2, 
                x='Category', 
                y='Value',
                color='Description',
                title="Electricity Generation (2000-2023)",
                labels={'Category': 'Energy Source'},
                height=600)
        
        fig11 = px.bar(df_pt_grouped, 
                x='Year', 
                y='EV/Revenue', 
                title="Recent Deals - EV/Revenue",
                labels={'EV/Revenue': 'EV/Revenue'},)
                # height=400)
        # fig11.update_traces(texttemplate='%{text:.1f}'+'x', textposition='auto',textfont=dict(size=10))

        fig12 = px.bar(df_pt_grouped, 
                x='Year', 
                y='EV/EBITDA', 
                title="Recent Deals - EV/EBITDA",
                labels={'EV/EBITDA': 'EV/EBITDA'},)
                # height=400)
        # fig12.update_traces(texttemplate='%{text:.1f}'+'x', textposition='auto',textfont=dict(size=10))

        fig13 = px.bar(
            df_rma_is,
            x='LineItems',
            y='Value',
            title="RMA - Income Statement",
            labels={'Value': 'Value ($)', 'LineItems': ' '},
            text_auto=True
        )

        fig14 = px.bar(
            df_rma_bs_grouped,
            x='Grouped_LineItems',
            y='Value',
            color='ReportID',
            title="RMA - Assets and Liabilities & Equity",
            labels={'Value_in_$': 'Value ($)', 'Grouped_LineItems': ' '},
            text_auto=True
        )

        fig14.update_layout(
            xaxis_title=" ",
            yaxis_title="Value ($)",
            xaxis_tickangle=45
        )

        fig15 = px.bar(filter_per_cap_electricity, 
                y='country', 
                x=['fossil', 'nuclear', 'renewable'], 
                title="Energy Mix per Capita by Country in 2023",
                labels={'value': 'Percentage', 'variable': 'Energy Source', 'country': 'Country'},
                color_discrete_map={'fossil': PRIMARY_COLORS['dark_blue'], 'nuclear': PRIMARY_COLORS['orange'], 'renewable': PRIMARY_COLORS['turquoise_blue']},
                orientation='h', 
                text_auto='.1%')  

        fig15.update_layout(
            xaxis=dict(title='Percentage', tickformat='.0%', range=[0, 1]),
            yaxis_title="Country",
            barmode='stack',
            legend_title="Energy Source",
            # height=500,
            # width=800
        )

        # Prices to Ultimate Customers
        fig16 = px.line(
            price_customers,
            x='Year',
            y='Value',
            color='Prices to Ultimate Customers',
            labels={
                'Value': 'Value',
                'Year': 'Year',
                'Prices to Ultimate Customers': 'Prices to Ultimate Customers'
            },
            title='Energy Prices',
            color_discrete_map=PRIMARY_COLORS
        )

        # Update layout for better appearance
        fig16.update_layout(
            xaxis_title='',
            yaxis_title='Value (cent per kilowatthour)',
            legend_title='',
            plot_bgcolor='rgba(0,0,0,0)',
            bargap=0.15,
            bargroupgap=0.1,
            legend=dict(
                orientation='h',
                yanchor='top',
                y=1.1,
                xanchor='left',
                x=0
            )
        )

        # Sales to Ultimate Customers
        fig17 = px.bar(
            sales_customers1,
            x='Year',
            y='Value',
            color='Sales to Ultimate Customers',
            labels={
                'Value': 'Value',
                'Year': 'Year',
                'Sales to Ultimate Customers': 'Sales to Ultimate Customers'
            },
            title='Energy Sales',
            barmode='stack',
            color_discrete_map=PRIMARY_COLORS
        )

        # Update layout for better appearance
        fig17.update_layout(
            xaxis_title='',
            yaxis_title='Value (billion kilowatthours)',
            legend_title='',
            plot_bgcolor='rgba(0,0,0,0)',
            bargap=0.15,
            bargroupgap=0.1,
            legend=dict(
                orientation='h',
                yanchor='top',
                y=1.1,
                xanchor='left',
                x=0
            )
        )

    # col1, col2 = st.columns(2)

    # with col1:

    st.plotly_chart(fig1, use_container_width=True)
    st.write("<h3 style='font-weight: bold; font-size:24px;'>Value Chain</h3>", unsafe_allow_html=True)
    st.image("https://www.energy-uk.org.uk/wp-content/uploads/2023/04/EUK-Different-parts-of-energy-market-diagram.webp", use_container_width=False)
    st.plotly_chart(fig2, use_container_width=True)
    # st.plotly_chart(fig3, use_container_width=True)
    st.plotly_chart(fig7, use_container_width=True)
        
    # with col2:
    st.plotly_chart(fig4, use_container_width=True)
    st.plotly_chart(fig5, use_container_width=True)
    st.plotly_chart(fig6, use_container_width=True)
    st.plotly_chart(fig8, use_container_width=True)
    
    st.plotly_chart(fig10, use_container_width=True)
    st.plotly_chart(fig11, use_container_width=True)
    st.plotly_chart(fig12, use_container_width=True)
    st.plotly_chart(fig13, use_container_width=True)
    st.plotly_chart(fig14, use_container_width=True)
    st.plotly_chart(fig15, use_container_width=True)
    st.plotly_chart(fig16, use_container_width=True)
    st.plotly_chart(fig17, use_container_width=True)

st.markdown("<h2 style='font-weight: bold; font-size:24px;'>Agriculture</h2>", unsafe_allow_html=True)
with st.expander("", expanded=False): 
    st.write("Agriculture-related analysis and visualizations go here.")

st.markdown("<h2 style='font-weight: bold; font-size:24px;'>Technology</h2>", unsafe_allow_html=True)
with st.expander("", expanded=False):
    st.write("Technology-related analysis and visualizations go here.")

st.markdown("<h2 style='font-weight: bold; font-size:24px;'>Automobiles</h2>", unsafe_allow_html=True)
with st.expander("", expanded=False):
    st.write("Automobiles-related analysis and visualizations go here.")

def export_to_pptx(fig1, fig2, fig3, fig4, fig5, fig6, fig7, fig8, fig10, fig11, fig12, fig13, fig14, fig15, fig16, fig17, value_chain_image_path, solar_image_path):

    template_path = os.path.join(os.getcwd(), "streamlit_dashboard", "data","energy_template.pptx")
    prs = Presentation(template_path)

    def add_chart_to_slide(prs, slide_index, fig, left, top, width, height):
        slide = prs.slides[slide_index]
        img_stream = BytesIO()
        fig.write_image(img_stream, format="png", engine="kaleido")
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, left, top, width=width, height=height)

    def add_image_to_slide(prs, slide_index, image_path, left, top, width, height):
        slide = prs.slides[slide_index]
        if os.path.exists(image_path):
            img_stream = BytesIO()
            with open(image_path, "rb") as img_file:
                img_stream.write(img_file.read())
            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, left, top, width=width, height=height)
        else:
            raise FileNotFoundError(f"The image at {image_path} was not found. Please check the path.")

    chart_configurations = [
        (4, fig1, Inches(1.35), Inches(6.75), Inches(2.5), Inches(6.20)),  # Slide 1: Market Size
        (4, fig2, Inches(4), Inches(6.75), Inches(3), Inches(3.4)),  # Slide 2: Electricity End Use
        (4, value_chain_image_path, Inches(0.3), Inches(4), Inches(2.5), Inches(6.20)),  # Slide 3: Value Chain
        (5, solar_image_path, Inches(1.3), Inches(0.3), Inches(2.5), Inches(6.25)),  # Slide 4: Solar
        (26, fig16, Inches(5.65), Inches(2.5), Inches(7), Inches(4.5)),  # Slide 26: Prices to Ultimate Customers 
        (26, fig17, Inches(5.65), Inches(2.5), Inches(7), Inches(4.5)),  # Slide 26: Sales to Ultimate Customers 
        (4, fig4, Inches(6.40), Inches(2.65), Inches(0.25), Inches(1.3)),  # Slide 5: Electricity Generation
        (5, fig5, Inches(1.3), Inches(6.70), Inches(2.5), Inches(6.25)),  # Slide 6: Renewable Share
        (6, fig6, Inches(1), Inches(1), Inches(8), Inches(5)),  # Slide 7: Per Capita Electricity
        (7, fig7, Inches(1), Inches(1), Inches(8), Inches(5)),  # Slide 8: Energy Source Consumption
        (8, fig8, Inches(1), Inches(1), Inches(8), Inches(5)),  # Slide 9: Energy Source Distribution
        (10, fig10, Inches(0.2), Inches(4.5), Inches(6.45), Inches(2.4)),  # Slide 11: Per Capita Generation
        (33, fig11, Inches(1.2), Inches(0.65), Inches(11), Inches(3.75)),  # Slide 12: Precedent Transaction - EV/Revenue
        (33, fig12, Inches(1.2), Inches(3.5), Inches(11), Inches(3.75)),  # Slide 13: Precedent Transaction - EV/EBITDA
        (36, fig13, Inches(1), Inches(1), Inches(8), Inches(5)),  # Slide 14: RMA - Income Statement
        (36, fig14, Inches(1), Inches(1), Inches(8), Inches(5)),  # Slide 15: RMA - Balance Sheet
        (6, fig15, Inches(5.65), Inches(2.5), Inches(7), Inches(4.5)),  # Slide 6: Per Capita Electricity Generation From Various Sources By Countries (%) 
    ]

    for config in chart_configurations:
        slide_index, content, left, top, width, height = config
        if isinstance(content, str):
            add_image_to_slide(prs, slide_index, content, left, top, width, height)
        else: 
            add_chart_to_slide(prs, slide_index, content, left, top, width, height)

    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream

def export_chart_options(fig1, fig2, fig4, fig5, fig6, fig7, fig8, fig10, fig11, fig12, fig13, fig14, fig15, fig16, fig17, value_chain_image_path, solar_image_path):
    if st.button("Export Charts to PowerPoint"):
        try:
            pptx_file = export_to_pptx(fig1, fig2, fig4, fig5, fig6, fig7, fig8, fig10, fig11, fig12, fig13, fig14, fig15, fig16, fig17, value_chain_image_path, solar_image_path)
            st.download_button(
                label="Download PowerPoint",
                data=pptx_file,
                file_name="Energy_Industry_Analysis_Report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except FileNotFoundError as e:
            st.error(f"Error: {e}")

value_chain_image_path = r"/mount/src/mck_cfo/streamlit_dashboard/data/value_chain.png"
solar_image_path = r"/mount/src/mck_cfo/streamlit_dashboard/data/solar.png"
export_chart_options(fig1, fig2, fig4, fig5, fig6, fig7, fig8, fig10, fig11, fig12, fig13, fig14, fig15, fig16, fig17, value_chain_image_path, solar_image_path)