import pandas as pd
import tensorflow as tf
import mysql.connector
import streamlit as st
import plotly.express as px
from st_aggrid import AgGrid, GridOptionsBuilder, GridUpdateMode
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import os
import threading

st.set_page_config(page_title="Precedent Transactions", layout="wide")

# MySQL database connection details
host = st.secrets["mysql"]["host"]
user = st.secrets["mysql"]["user"]
password = st.secrets["mysql"]["password"]
database = st.secrets["mysql"]["database"]

# Connect to the MySQL database
try:
    conn = mysql.connector.connect(
        host=host,
        user=user,
        password=password,
        database=database
    )
except mysql.connector.Error as e:
    st.error(f"Error connecting to MySQL: {e}")
    st.stop()

@st.cache_data(ttl=600)  # Cache for 10 minutes
def load_data():
    conn = mysql.connector.connect(
        host=host,
        user=user,
        password=password,
        database=database
    )
    query = """
    SELECT `Year`, `Target`, `EV/Revenue`, `EV/EBITDA`, `Business Description`, `Industry`, `Location`
    FROM precedent_transaction
    WHERE Industry IS NOT NULL AND Location IS NOT NULL
    """
    df = pd.read_sql_query(query, conn)
    conn.close()
    return df

df = load_data()

# Ensure all columns are TensorFlow-compatible
df['Year'] = df['Year'].fillna(0).astype('int32')
df['EV/Revenue'] = df['EV/Revenue'].fillna(0).astype('float32').round(2)
df['EV/EBITDA'] = df['EV/EBITDA'].fillna(0).astype('float32').round(2)
df['Target'] = df['Target'].fillna("").astype('string')
df['Business Description'] = df['Business Description'].fillna("").astype('string')
df['Industry'] = df['Industry'].fillna("").astype('string')
df['Location'] = df['Location'].fillna("").astype('string')

# List of unique industries and locations
industries = df['Industry'].unique()
locations = df['Location'].unique()

# Sidebar for selecting industries and locations
col1, col2 = st.columns(2)
selected_industries = col1.multiselect("Select Industry", industries)
selected_locations = col2.multiselect("Select Location", locations)

if selected_industries and selected_locations:

    filtered_precedent_df = df[df['Industry'].isin(selected_industries) & df['Location'].isin(selected_locations)]
    filtered_precedent_df = filtered_precedent_df[['Target', 'Year', 'EV/Revenue', 'EV/EBITDA', 'Business Description']]
    filtered_precedent_df['Year'] = filtered_precedent_df['Year'].astype(int)
    filtered_precedent_df['EV/Revenue'] = filtered_precedent_df['EV/Revenue'].astype('float32').round(2)
    filtered_precedent_df['EV/EBITDA'] = filtered_precedent_df['EV/EBITDA'].astype('float32').round(2)

    # Display filtered data in Ag-Grid table
    st.subheader("Precedent Transactions")
    gb = GridOptionsBuilder.from_dataframe(filtered_precedent_df)
    gb.configure_selection(selection_mode="multiple", use_checkbox=True)
    gb.configure_column(
        field="Target",
        tooltipField="Business Description",
        maxWidth=400
    )
    gb.configure_columns(["Business Description"], hide=False)
    grid_options = gb.build()

    # Display Ag-Grid table
    grid_response = AgGrid(
        filtered_precedent_df,
        gridOptions=grid_options,
        update_mode=GridUpdateMode.SELECTION_CHANGED,
        height=400,
        width='100%',
        theme='streamlit'
    )

    # Get selected rows from the grid
    selected_data = pd.DataFrame(grid_response['selected_rows'])
    if not selected_data.empty:
        # Calculate average EV/Revenue and EV/EBITDA for the selected data
        avg_data = selected_data.groupby('Year')[['EV/Revenue', 'EV/EBITDA']].mean().reset_index()
        avg_data['Year'] = avg_data['Year'].astype(int)
        avg_data['EV/EBITDA'] = avg_data['EV/EBITDA'].round(2)
        avg_data['EV/Revenue'] = avg_data['EV/Revenue'].round(2)
        color_ev_revenue = "#032649" 
        color_ev_ebitda = "#032649"  

        median_ev_revenue = avg_data['EV/Revenue'].median()
        median_ev_ebitda = avg_data['EV/EBITDA'].median()


        fig1_precedent = px.bar(avg_data, x='Year', y='EV/Revenue', title="EV/Revenue", text='EV/Revenue')  # No title
        fig1_precedent.update_traces(marker_color=color_ev_revenue, texttemplate='%{text:.1f}'+'x', textposition='auto',textfont=dict(size=12))
        fig1_precedent.update_layout(yaxis_title="EV/Revenue", xaxis_title=" ", bargap=0.4, bargroupgap=0.4, yaxis=dict(showgrid=False),xaxis=dict(tickmode='linear', tick0=avg_data['Year'].min(), dtick=1), shapes=[dict(type='line', x0=avg_data['Year'].min(), x1=avg_data['Year'].max(), y0=median_ev_revenue, y1=median_ev_revenue, line=dict(color='#EB8928', dash='dot', width=2))], annotations=[dict(x=avg_data['Year'].max(), y=median_ev_revenue, xanchor='left', yanchor='bottom', text=f'Median: {median_ev_revenue:.1f}'+'x', showarrow=False, font=dict(size=12, color='gray'), bgcolor='white')],plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',margin=dict(l=0, r=0, t=0),width=900,height=300)

        st.plotly_chart(fig1_precedent)

        fig2_precedent = px.bar(avg_data, x='Year', y='EV/EBITDA', title="EV/EBITDA", text='EV/EBITDA')
        fig2_precedent.update_traces(marker_color=color_ev_ebitda, texttemplate='%{text:.1f}'+ 'x', textposition='auto',textfont=dict(size=12))
        fig2_precedent.update_layout(yaxis_title="EV/EBITDA", xaxis_title=" ", bargap=0.4, bargroupgap=0.4, yaxis=dict(showgrid=False),xaxis=dict(tickmode='linear', tick0=avg_data['Year'].min(), dtick=1), shapes=[dict(type='line', x0=avg_data['Year'].min(), x1=avg_data['Year'].max(), y0=median_ev_ebitda, y1=median_ev_ebitda, line=dict(color='#EB8928', dash='dot', width=2))], annotations=[dict(x=avg_data['Year'].max(), y=median_ev_ebitda, xanchor='left', yanchor='bottom', text=f'Median: {median_ev_ebitda:.1f}'+'x', showarrow=False, font=dict(size=12, color='gray'), bgcolor='white')],plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',margin=dict(l=0, r=0, t=0),width=900,height=300)
            
        st.plotly_chart(fig2_precedent)

        export_ppt = st.button("Export Charts to PowerPoint")

        def generate_ppt():
            template_path = os.path.join(os.getcwd(), "streamlit_dashboard", "data", "main_template_pitch.pptx")
            if not os.path.exists(template_path):
                st.error("PowerPoint template not found!")
                return None

            ppt = Presentation(template_path)
            slide1 = ppt.slides[10]

            fig1_precedent_image = BytesIO()
            fig1_precedent.write_image(fig1_precedent_image, format="png", width=900, height=300)
            fig1_precedent_image.seek(0)
            slide1.shapes.add_picture(fig1_precedent_image, Inches(0.11), Inches(0.90), width=Inches(9), height=Inches(2.8))

            fig2_precedent_image = BytesIO()
            fig2_precedent.write_image(fig2_precedent_image, format="png", width=900, height=300)
            fig2_precedent_image.seek(0)
            slide1.shapes.add_picture(fig2_precedent_image, Inches(0.11), Inches(3.70), width=Inches(9), height=Inches(2.8))

            ppt_bytes = BytesIO()
            ppt.save(ppt_bytes)
            ppt_bytes.seek(0)

            return ppt_bytes

        if export_ppt:
            ppt_thread = threading.Thread(target=generate_ppt)
            ppt_thread.start()

            ppt_bytes = generate_ppt()
            if ppt_bytes:
                st.download_button(
                    label="Download PowerPoint",
                    data=ppt_bytes,
                    file_name="precedent_transaction.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

else:
    st.warning("Please select at least one Industry and Location to view data.")