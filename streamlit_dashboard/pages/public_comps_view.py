# pages/public_comps_view.py
import os
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from st_aggrid import AgGrid, GridOptionsBuilder, GridUpdateMode

# Path to the data
# user_profile = os.path.expanduser("~")
# source_path = os.path.join(user_profile, 'source', 'mck_setup', 'streamlit_dashboard', 'data')
path_public_comps= os.path.join(
    os.getcwd(), "streamlit_dashboard", "data", "Updated - Public Listed Companies US.xlsx"
)

def get_public_comps_layout():
    """Render the Public Companies page layout."""
    st.header("Public Companies")

    # Load and process data
    public_comps_df = pd.read_excel(path_public_comps, sheet_name="FY 2023")
    public_comps_df['Enterprise Value (in $)'] = pd.to_numeric(public_comps_df['Enterprise Value (in $)'], errors='coerce')
    public_comps_df['Revenue (in $)'] = pd.to_numeric(public_comps_df['Revenue (in $)'], errors='coerce').round(1)
    public_comps_df['EBITDA (in $)'] = pd.to_numeric(public_comps_df['EBITDA (in $)'], errors='coerce').round(1)
    public_comps_df['EV/Revenue'] = public_comps_df['Enterprise Value (in $)'] / public_comps_df['Revenue (in $)']
    public_comps_df['EV/EBITDA'] = public_comps_df['Enterprise Value (in $)'] / public_comps_df['EBITDA (in $)']

    columns_to_display = ['Name', 'Country', 'Industry', 'EV/Revenue', 'EV/EBITDA', 'Business Description']
    filtered_df = public_comps_df[columns_to_display]

    # Configure AgGrid
    gb = GridOptionsBuilder.from_dataframe(filtered_df)
    gb.configure_selection('multiple', use_checkbox=True)
    gb.configure_default_column(editable=True, filter=True, sortable=True, resizable=True)
    grid_options = gb.build()

    # st.subheader("Company Data")
    grid_response = AgGrid(
        filtered_df,
        gridOptions=grid_options,
        update_mode=GridUpdateMode.SELECTION_CHANGED,
        theme='streamlit',
        fit_columns_on_grid_load=True,
        height=500,
        width='100%'
    )

    selected_rows = pd.DataFrame(grid_response['selected_rows'])

    if not selected_rows.empty:
        ev_revenue_fig, ev_ebitda_fig = plot_public_comps_charts(selected_rows)
        export_chart_options(ev_revenue_fig, ev_ebitda_fig)
    else:
        st.info("Select companies to visualize their data.")
        
sns.set_style("whitegrid")
plt.rc('axes', edgecolor='gray')
plt.rc('xtick', color='gray')
plt.rc('ytick', color='gray')

def plot_public_comps_charts(data):
    """Plot EV/Revenue and EV/EBITDA charts."""
    # EV/Revenue Chart
    fig1, ax1 = plt.subplots(figsize=(12, 4))
    sns.barplot(data=data, x='Name', y='EV/Revenue', ax=ax1, color='#032649')
    for p in ax1.patches:
        ax1.annotate(f'{p.get_height():.1f}', (p.get_x() + p.get_width() / 2., p.get_height()),
                     ha='center', va='center', fontsize=10, color='black', xytext=(0, 5), textcoords='offset points')
    st.pyplot(fig1)

    # EV/EBITDA Chart
    fig2, ax2 = plt.subplots(figsize=(12, 4))
    sns.barplot(data=data, x='Name', y='EV/EBITDA', ax=ax2, color='#EB8928')
    for p in ax2.patches:
        ax2.annotate(f'{p.get_height():.1f}', (p.get_x() + p.get_width() / 2., p.get_height()),
                     ha='center', va='center', fontsize=10, color='black', xytext=(0, 5), textcoords='offset points')
    st.pyplot(fig2)

    return fig1, fig2

def export_chart_options(ev_revenue_fig, ev_ebitda_fig):
    """Export charts as PNG or PowerPoint."""
    st.subheader("Export Charts")

    # Download as PNG
    png_buffer1 = BytesIO()
    ev_revenue_fig.savefig(png_buffer1, format="png")
    png_buffer1.seek(0)

    png_buffer2 = BytesIO()
    ev_ebitda_fig.savefig(png_buffer2, format="png")
    png_buffer2.seek(0)

    # st.download_button(
    #     label="Download EV/Revenue Chart as PNG",
    #     data=png_buffer1,
    #     file_name="ev_revenue_chart.png",
    #     mime="image/png"
    # )

    # st.download_button(
    #     label="Download EV/EBITDA Chart as PNG",
    #     data=png_buffer2,
    #     file_name="ev_ebitda_chart.png",
    #     mime="image/png"
    # )

    # Export to PowerPoint
    if st.button("Export Charts to PowerPoint"):
        pptx_file = export_to_pptx(ev_revenue_fig, ev_ebitda_fig)
        st.download_button(
            label="Download PowerPoint",
            data=pptx_file,
            file_name="public_companies_charts.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

def export_to_pptx(ev_revenue_fig, ev_ebitda_fig):
    """Export charts to a PowerPoint presentation."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]

    slide1 = prs.slides.add_slide(slide_layout)
    title1 = slide1.shapes.title
    title1.text = "EV/Revenue Chart"
    img1 = BytesIO()
    ev_revenue_fig.savefig(img1, format="png", bbox_inches='tight')
    img1.seek(0)
    slide1.shapes.add_picture(img1, Inches(1), Inches(1), width=Inches(10), height=Inches(4.5))

    slide2 = prs.slides.add_slide(slide_layout)
    title2 = slide2.shapes.title
    title2.text = "EV/EBITDA Chart"
    img2 = BytesIO()
    ev_ebitda_fig.savefig(img2, format="png", bbox_inches='tight')
    img2.seek(0)
    slide2.shapes.add_picture(img2, Inches(1), Inches(1), width=Inches(10), height=Inches(4.5))

    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io
