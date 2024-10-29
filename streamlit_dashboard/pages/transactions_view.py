import os
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from st_aggrid import AgGrid, GridOptionsBuilder, GridUpdateMode

path_transaction = os.path.abspath(r'streamlit_dashboard/data/Updated - Precedent Transaction.xlsx')
@st.cache_data
def get_transaction_layout():    
    transactions_df = pd.read_excel(path_transaction, sheet_name="Final - Precedent Transactions")
    transactions_df['Announced Date'] = pd.to_datetime(transactions_df['Announced Date'], errors='coerce')
    transactions_df.dropna(subset=['Announced Date'], inplace=True)
    transactions_df['Year'] = transactions_df['Announced Date'].dt.year.astype(int)
    transactions_df['EV/Revenue'] = pd.to_numeric(transactions_df['EV/Revenue'], errors='coerce').fillna(0).round(1)
    transactions_df['EV/EBITDA'] = pd.to_numeric(transactions_df['EV/EBITDA'], errors='coerce').fillna(0).round(1)
    columns_to_display = {
        'Target': 'Company',
        'Geographic Locations': 'Location',
        'Year': 'Year',
        'Industry': 'Industry',
        'EV/Revenue': 'EV/Revenue',
        'EV/EBITDA': 'EV/EBITDA',
        'Business Description': 'Business Description'
    }
    filtered_df = transactions_df[list(columns_to_display.keys())].rename(columns=columns_to_display)

    gb = GridOptionsBuilder.from_dataframe(filtered_df)
    gb.configure_selection('multiple', use_checkbox=True)
    gb.configure_default_column(editable=True, filter=True, sortable=True, resizable=True)
    grid_options = gb.build()
    grid_response = AgGrid(
        filtered_df,
        gridOptions=grid_options,
        update_mode=GridUpdateMode.SELECTION_CHANGED,
        theme='alpine',
        fit_columns_on_grid_load=True,
        height=500,
        width='100%'
    )

    selected_rows = pd.DataFrame(grid_response['selected_rows'])

    bar_width = st.sidebar.slider("Select Bar Width", min_value=0.1, max_value=0.9, value=0.5, step=0.1)

    if not selected_rows.empty:
        ev_revenue_fig, ev_ebitda_fig = plot_transactions_charts(selected_rows, bar_width)
        export_chart_options(ev_revenue_fig, ev_ebitda_fig)
    else:
        st.info("Select companies to visualize their data.")
        
sns.set_style("whitegrid")
plt.rc('axes', edgecolor='gray')
plt.rc('xtick', color='gray')
plt.rc('ytick', color='gray')

def plot_transactions_charts(data, bar_width):
    """Plot EV/Revenue and EV/EBITDA charts."""
    grouped_data = data.groupby('Year').agg(
        avg_ev_revenue=('EV/Revenue', 'mean'),
        avg_ev_ebitda=('EV/EBITDA', 'mean')
    ).reset_index()

    # EV/Revenue Chart
    st.subheader("EV/Revenue")
    fig1, ax1 = plt.subplots(figsize=(12, 4))
    sns.barplot(data=grouped_data, x='Year', y='avg_ev_revenue', ax=ax1, color='#032649', width=bar_width)

    for p in ax1.patches:
        ax1.annotate(f'{p.get_height():.1f}x',
                     (p.get_x() + p.get_width() / 2., p.get_height()),
                     ha='center', va='center', fontsize=10, color='black',
                     xytext=(0, 5), textcoords='offset points')

    ax1.set_xlabel("")  # Remove x-axis label
    ax1.set_xticklabels(ax1.get_xticklabels(), rotation=45, ha='right')
    st.pyplot(fig1)
    plt.close(fig1)

    # EV/EBITDA Chart
    st.subheader("EV/EBITDA")
    fig2, ax2 = plt.subplots(figsize=(12, 4))
    sns.barplot(data=grouped_data, x='Year', y='avg_ev_ebitda', ax=ax2, color='#EB8928', width=bar_width)

    for p in ax2.patches:
        ax2.annotate(f'{p.get_height():.1f}x',
                     (p.get_x() + p.get_width() / 2., p.get_height()),
                     ha='center', va='center', fontsize=10, color='black',
                     xytext=(0, 5), textcoords='offset points')

    ax2.set_xlabel("")  # Remove x-axis label
    ax2.set_xticklabels(ax2.get_xticklabels(), rotation=45, ha='right')
    st.pyplot(fig2)
    plt.close(fig2)

    return fig1, fig2

def export_chart_options(ev_revenue_fig, ev_ebitda_fig):
    """Export charts as PNG or PowerPoint."""
    st.subheader("Export Charts")

    # Download as PNG
    png_buffer1 = BytesIO()
    ev_revenue_fig.savefig(png_buffer1, format="png")
    png_buffer1.seek(0)

    png_buffer2 = BytesIO()
    ev_ebitda_fig.savefig(png_buffer2, format="png")
    png_buffer2.seek(0)

    # st.download_button(
    #     label="Download EV/Revenue Chart as PNG",
    #     data=png_buffer1,
    #     file_name="ev_revenue_chart.png",
    #     mime="image/png"
    # )

    # st.download_button(
    #     label="Download EV/EBITDA Chart as PNG",
    #     data=png_buffer2,
    #     file_name="ev_ebitda_chart.png",
    #     mime="image/png"
    # )

    # Export to PowerPoint
    if st.button("Export Charts to PowerPoint"):
        pptx_file = export_to_pptx(ev_revenue_fig, ev_ebitda_fig)
        st.download_button(
            label="Download PowerPoint",
            data=pptx_file,
            file_name="public_companies_charts.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

def export_to_pptx(ev_revenue_fig, ev_ebitda_fig):
    """Export charts to a PowerPoint presentation."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]

    slide1 = prs.slides.add_slide(slide_layout)
    title1 = slide1.shapes.title
    title1.text = "EV/Revenue Chart"
    img1 = BytesIO()
    ev_revenue_fig.savefig(img1, format="png", bbox_inches='tight')
    img1.seek(0)
    slide1.shapes.add_picture(img1, Inches(0.5), Inches(1.15), width=Inches(9), height=Inches(2.8))

    # slide2 = prs.slides.add_slide(slide_layout)
    # title2 = slide2.shapes.title
    # title2.text = "EV/EBITDA Chart"
    img2 = BytesIO()
    ev_ebitda_fig.savefig(img2, format="png", bbox_inches='tight')
    img2.seek(0)
    slide1.shapes.add_picture(img2, Inches(0.5), Inches(4.35), width=Inches(9), height=Inches(2.8))

    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io
