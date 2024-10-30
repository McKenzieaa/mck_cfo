import streamlit as st
import pandas as pd
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
import time
from io import BytesIO
import plotly.io as pio

# Set Kaleido renderer for image export
pio.kaleido.scope.default_format = "png"

# Import chart functions from other modules
from pages.public_comps_view import plot_public_comps_charts
from pages.transactions_view import plot_transactions_charts
from pages.us_state_indicators_view import plot_gdp_chart,plot_unemployment_labour_chart

def presentation_view():
    """Generate a presentation dynamically from selected charts."""
    st.title("Generate Custom Presentation from Charts")

    # Store available chart functions
    chart_functions = {
        "Public Comps": plot_public_comps_charts,
        "Precedent Transactions": plot_transactions_charts,
        "State Indicators": (plot_gdp_chart,plot_unemployment_labour_chart)
    }

    # User interface to select charts
    selected_charts = st.multiselect(
        "Select Charts to Include in Presentation", 
        options=list(chart_functions.keys())
    )

    # Button to generate the presentation
    if st.button("Generate Presentation"):
        if not selected_charts:
            st.warning("Please select at least one chart.")
        else:
            with st.spinner("Generating presentation..."):
                start_time = time.time()

                # Create a new PowerPoint presentation
                prs = Presentation()

                # Loop through selected charts and add them to the presentation
                for chart_name in selected_charts:
                    if chart_name in chart_functions:
                        chart_func = chart_functions[chart_name]
                        fig = chart_func()  # Generate the chart

                        # Validate the generated chart
                        if fig is None or not isinstance(fig, go.Figure):
                            st.error(f"Failed to generate chart: {chart_name}")
                            continue  # Skip this chart if it fails

                        # Save chart as image in memory
                        img_stream = BytesIO()
                        try:
                            fig.write_image(img_stream, format="png")
                        except Exception as e:
                            st.error(f"Error generating chart image: {str(e)}")
                            continue  # Skip this chart if image export fails

                        img_stream.seek(0)  # Reset the stream position

                        # Add slide and image to presentation
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        slide.shapes.title.text = chart_name
                        slide.shapes.add_picture(
                            img_stream, Inches(1), Inches(1.5), 
                            width=Inches(8), height=Inches(4.5)
                        )

                # Save presentation in memory
                presentation_stream = BytesIO()
                prs.save(presentation_stream)
                presentation_stream.seek(0)

                # Total time taken
                total_time = time.time() - start_time
                st.success(f"Presentation generated in {total_time:.2f} seconds.")

                # Provide download button
                st.download_button(
                    label="Download Presentation",
                    data=presentation_stream,
                    file_name="custom_presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

presentation_view()