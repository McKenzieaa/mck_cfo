import streamlit as st
import pandas as pd
import plotly.graph_objs as go
from sqlalchemy import create_engine
from datetime import date
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# MySQL Database Connection
host = st.secrets["mysql"]["host"]
user = st.secrets["mysql"]["user"]
password = st.secrets["mysql"]["password"]
database = st.secrets["mysql"]["database"]
engine = create_engine(f"mysql+pymysql://{user}:{password}@{host}/{database}")

# Function to fetch state economic data
def fetch_state_data(state_name):
    query = f"""
    SELECT year, month, unemployment_rate, labour_force, gdp, population
    FROM state_economic_data
    WHERE state = '{state_name}' AND year >= 2000
    ORDER BY year, month;
    """
    try:
        with engine.connect() as connection:
            df = pd.read_sql(query, connection)
        return df
    except Exception as e:
        st.error(f"Error fetching data: {e}")
        return pd.DataFrame()

# Define Plot Colors
line_colors = {
    "unemployment": "#032649",  # Dark Blue
    "labour_force": "#EB8928",  # Orange
    "gdp": "#032649",  # Dark Blue
    "population": "rgba(0, 102, 255, 0.4)",  # Light Transparent Blue
}

# Function to plot Unemployment & Labour Force Chart
def plot_unemployment_labour_chart(state_name):
    df = fetch_state_data(state_name)
    
    if df.empty:
        st.warning(f"No data available for {state_name}.")
        return None
    
    df["date"] = pd.to_datetime(df["year"].astype(str) + "-" + df["month"].astype(str) + "-01")

    fig = go.Figure()
    fig.add_trace(go.Scatter(x=df["date"], y=df["unemployment_rate"], mode='lines',
                             line=dict(color=line_colors["unemployment"]), name="Unemployment Rate"))
    fig.add_trace(go.Scatter(x=df["date"], y=df["labour_force"], mode='lines',
                             line=dict(color=line_colors["labour_force"]), name="Labour Force"))

    fig.update_layout(
        title="",
        xaxis_title=" ",
        yaxis_title="Rate",
        template="plotly_white",
        legend=dict(
            x=0.01, 
            y=-0.2,  
            xanchor='left',
            yanchor='top',
            title_text=None,
            orientation='h', 
            font=dict(size=10)
        ),
        plot_bgcolor='rgba(0,0,0,0)',
        paper_bgcolor='rgba(0,0,0,0)',
        margin=dict(l=5, r=5, t=10, b=120), 
        height=300,
        width=500,
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False)
    )

    st.plotly_chart(fig, use_container_width=True)
    return fig

# Function to plot GDP Chart
def plot_gdp_chart(state_name):
    df = fetch_state_data(state_name)
    
    if df.empty:
        st.warning(f"No GDP data available for {state_name}.")
        return None
    
    df = df.groupby("year").agg({"gdp": "sum"}).reset_index()

    fig = go.Figure()
    fig.add_trace(go.Scatter(x=df['year'], y=df['gdp'], mode='lines',
                             line=dict(color=line_colors["gdp"]), name=f"{state_name} GDP"))

    fig.update_layout(
        title="",
        xaxis_title=" ",
        yaxis_title="State GDP",
        template="plotly_white",
        legend=dict(
            x=0.01, 
            y=0.01, 
            xanchor='left', 
            yanchor='bottom', 
            title_text=None
        ),
        plot_bgcolor='rgba(0,0,0,0)',
        paper_bgcolor='rgba(0,0,0,0)',
        margin=dict(l=2, r=2, t=30, b=50),
        height=300,
        width=500,
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False)
    )

    st.plotly_chart(fig, use_container_width=True)
    return fig

# Function to plot Population Area Chart
def plot_population_chart(state_name):
    df = fetch_state_data(state_name)

    if df.empty:
        st.warning(f"No population data available for {state_name}.")
        return None

    df = df.groupby("year").agg({"population": "max"}).reset_index()

    fig = go.Figure()
    fig.add_trace(go.Scatter(
        x=df["year"],
        y=df["population"],
        fill='tozeroy',  # Area fill
        mode='lines',
        line=dict(color=line_colors["population"]),
        name="Population"
    ))

    fig.update_layout(
        title="",
        xaxis_title=" ",
        yaxis_title="Population",
        template="plotly_white",
        legend=dict(
            x=0.01, 
            y=0.01, 
            xanchor='left', 
            yanchor='bottom', 
            title_text=None
        ),
        plot_bgcolor='rgba(0,0,0,0)',
        paper_bgcolor='rgba(0,0,0,0)',
        margin=dict(l=2, r=2, t=30, b=50),
        height=300,
        width=500,
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False)
    )

    st.plotly_chart(fig, use_container_width=True)
    return fig

# Function to export charts to PowerPoint
def export_to_pptx(labour_fig, gdp_fig, population_fig):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "State Economic Indicators"

    img1 = BytesIO()
    labour_fig.write_image(img1, format="png")
    img1.seek(0)
    slide.shapes.add_picture(img1, Inches(0.5), Inches(0.5), width=Inches(4.5), height=Inches(3.5))

    img2 = BytesIO()
    gdp_fig.write_image(img2, format="png")
    img2.seek(0)
    slide.shapes.add_picture(img2, Inches(5), Inches(0.5), width=Inches(4.5), height=Inches(3.5))

    img3 = BytesIO()
    population_fig.write_image(img3, format="png")
    img3.seek(0)
    slide.shapes.add_picture(img3, Inches(2.75), Inches(4), width=Inches(4.5), height=Inches(3.5))

    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# Streamlit App Layout
def get_state_indicators_layout():
    st.title("US State Indicators")

    query_states = "SELECT DISTINCT state FROM state_economic_data;"
    states = fetch_state_data("")["state"].unique().tolist() if fetch_state_data("") is not None else []
    
    if states:
        state_name = st.selectbox("Select State", states, index=0)

        st.subheader(f"{state_name} - Unemployment & Labour Force")
        labour_fig = plot_unemployment_labour_chart(state_name)

        st.subheader(f"{state_name} - GDP Over Time")
        gdp_fig = plot_gdp_chart(state_name)

        st.subheader(f"{state_name} - Population Growth")
        population_fig = plot_population_chart(state_name)

        if st.button("Export Charts to PowerPoint") and labour_fig and gdp_fig and population_fig:
            pptx_file = export_to_pptx(labour_fig, gdp_fig, population_fig)
            st.download_button("Download PowerPoint", pptx_file, f"{state_name}_state_indicators.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    else:
        st.warning("No states found in the database.")

get_state_indicators_layout()
