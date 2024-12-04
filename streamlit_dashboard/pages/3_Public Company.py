import pandas as pd
from sqlalchemy import create_engine
import streamlit as st
import plotly.express as px
from st_aggrid import AgGrid, GridOptionsBuilder, GridUpdateMode
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import os

# Streamlit app title
st.set_page_config(page_title="Public Listed Companies Analysis", layout="wide")

# MySQL connection setup
mysql_user = st.secrets["mysql"]["user"]
mysql_password = st.secrets["mysql"]["password"]
mysql_host = st.secrets["mysql"]["host"]
mysql_db = st.secrets["mysql"]["db"]

# Create SQLAlchemy engine
connection_string = f"mysql+mysqlconnector://{mysql_user}:{mysql_password}@{mysql_host}/{mysql_db}"
engine = create_engine(connection_string)

# SQL query to fetch the public_comp_table
query = """
    SELECT `Name`, `Country`, `Enterprise Value (in $)`, `Revenue (in $)`, `EBITDA (in $)`, `Business Description`, `Industry`
    FROM public_comp_table
"""

# Load data from MySQL
try:
    df = pd.read_sql(query, engine)

    # Rename columns to match expected column names
    df = df.rename(columns={
        'Name': 'Company',
        'Country': 'Location',
        'Enterprise Value (in $)': 'Enterprise Value',
        'Revenue (in $)': 'Revenue',
        'EBITDA (in $)': 'EBITDA',
    })

    # Convert columns to numeric
    df['Enterprise Value'] = pd.to_numeric(df['Enterprise Value'], errors='coerce')
    df['Revenue'] = pd.to_numeric(df['Revenue'], errors='coerce')
    df['EBITDA'] = pd.to_numeric(df['EBITDA'], errors='coerce')

    # Calculate EV/Revenue and EV/EBITDA
    df['EV/Revenue'] = df['Enterprise Value'] / df['Revenue']
    df['EV/EBITDA'] = df['Enterprise Value'] / df['EBITDA']

except Exception as e:
    st.error(f"Error loading data from MySQL: {e}")
    st.stop()

# Get unique values for Industry and Location filters
industries = df['Industry'].dropna().unique()
locations = df['Location'].dropna().unique()

col1, col2 = st.columns(2)
selected_industries = col1.multiselect("Select Industry", industries)
selected_locations = col2.multiselect("Select Location", locations)

# Filter data based on multi-selections using .isin()
if selected_industries and selected_locations:
    filtered_df = df[df['Industry'].isin(selected_industries) & df['Location'].isin(selected_locations)]
    filtered_df = filtered_df[['Company',  'EV/Revenue', 'EV/EBITDA', 'Business Description']]
    filtered_df['EV/Revenue'] = filtered_df['EV/Revenue'].round(1)
    filtered_df['EV/EBITDA'] = filtered_df['EV/EBITDA'].round(1)

    # Set up Ag-Grid for selection
    st.title("Public Listed Companies")
    gb = GridOptionsBuilder.from_dataframe(filtered_df)
    gb.configure_selection(selection_mode="multiple", use_checkbox=True)
    gb.configure_column(
        field="Company",
        tooltipField="Business Description",
        maxWidth=400
    )
    gb.configure_columns(["Business Description"], hide=False)    
    grid_options = gb.build()

    # Display Ag-Grid table
    grid_response = AgGrid(
        filtered_df,
        gridOptions=grid_options,
        update_mode=GridUpdateMode.SELECTION_CHANGED,
        height=400,
        width='100%',
        theme='streamlit'
    )

    selected_data = pd.DataFrame(grid_response['selected_rows'])
    if not selected_data.empty:
        avg_data = selected_data.groupby('Company')[['EV/Revenue', 'EV/EBITDA']].mean().reset_index()
        avg_data['Company'] = avg_data['Company'].apply(lambda x: '<br>'.join([x[i:i+20] for i in range(0, len(x), 20)]) if len(x) > 20 else x)

        color_ev_revenue = "#032649"  # Default Plotly blue
        color_ev_ebitda = "#032649"   # Default Plotly red

        median_ev_revenue = avg_data['EV/Revenue'].median()
        median_ev_ebitda = avg_data['EV/EBITDA'].median()

        fig1_public = px.bar(avg_data, x='Company', y='EV/Revenue', title="EV/Revenue", text='EV/Revenue')
        fig1_public.update_traces(marker_color=color_ev_revenue, texttemplate='%{text:.1f}'+'x', textposition='auto',textfont=dict(size=12))
        fig1_public.update_layout(yaxis_title="EV/Revenue", xaxis_title=" ",bargap=0.4,bargroupgap=0.4,yaxis=dict(showgrid=False),xaxis=dict(tickangle=0,automargin=True,tickmode='array',tickvals=avg_data['Company'],ticktext=avg_data['Company']),plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',margin=dict(l=0, r=0, t=50,b=80),width=900,height=300)
        fig1_public.add_shape(type="line",x0=-0.5, x1=len(avg_data['Company']) - 0.5,  y0=median_ev_revenue, y1=median_ev_revenue,line=dict(color="#EB8928", width=2, dash="dot"),  xref="x", yref="y")
        fig1_public.add_annotation(x=len(avg_data['Company']) - 1, y=median_ev_revenue + 0.2, text=f"Median: {median_ev_revenue:.1f}x",showarrow=False, font=dict(size=10, color="gray"), xanchor="left",bgcolor='white')

        st.plotly_chart(fig1_public)

        fig2_public = px.bar(avg_data, x='Company', y='EV/EBITDA', title="EV/EBITDA", text='EV/EBITDA')
        fig2_public.update_traces(marker_color=color_ev_ebitda,texttemplate='%{text:.1f}'+'x', textposition='auto',textfont=dict(size=12))
        fig2_public.update_layout(yaxis_title="EV/EBITDA", xaxis_title=" ",bargap=0.4,bargroupgap=0.4,yaxis=dict(showgrid=False),xaxis=dict(tickangle=0,automargin=True,tickmode='array',tickvals=avg_data['Company'],ticktext=avg_data['Company']),plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',margin=dict(l=0, r=0, t=50,b=80),width=900,height=300)
        fig2_public.add_shape(type="line",x0=-0.5, x1=len(avg_data['Company']) - 0.5,  y0=median_ev_ebitda, y1=median_ev_ebitda,line=dict(color="#EB8928", width=2, dash="dot"),  xref="x", yref="y")
        fig2_public.add_annotation(x=len(avg_data['Company']) - 1, y=median_ev_ebitda + 0.2, text=f"Median: {median_ev_ebitda:.1f}x",showarrow=False, font=dict(size=10, color="gray"), xanchor="left",bgcolor='white')
            
        st.plotly_chart(fig2_public)
        export_ppt = st.button("Export Charts to PowerPoint")

        if export_ppt:

            template_path = os.path.join(os.getcwd(), "streamlit_dashboard", "data", "main_template_pitch.pptx")

            if not os.path.exists(template_path):
                st.error(f"PowerPoint template not found at: {template_path}")
                st.stop()

            ppt = Presentation(template_path)
            slide1 = ppt.slides[11] 

            if slide1 is None:
                slide_layout = ppt.slide_layouts[5]
                slide1 = ppt.slides.add_slide(slide_layout)

            title1 = slide1.shapes.title
            # title1.text = ""  # Remove chart title
            
            fig1_image = BytesIO()
            fig1_public.write_image(fig1_image, format="png", width=900, height=300)
            fig1_image.seek(0)
            slide1.shapes.add_picture(fig1_image, Inches(0.11), Inches(0.90), width=Inches(9), height=Inches(2.8))

            # Add EV/EBITDA chart to the same slide
            fig2_image = BytesIO()
            fig2_public.write_image(fig2_image, format="png", width=900, height=300)
            fig2_image.seek(0)
            slide1.shapes.add_picture(fig2_image, Inches(0.11), Inches(3.70), width=Inches(9), height=Inches(2.8))

            # Save PowerPoint to BytesIO object for download
            ppt_bytes = BytesIO()
            ppt.save(ppt_bytes)
            ppt_bytes.seek(0)

            # Provide download link for PowerPoint
            st.download_button(
                label="Download PowerPoint",
                data=ppt_bytes,
                file_name="public_comps.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

else:
    st.write("Please select at least one Industry and Location to view data.")